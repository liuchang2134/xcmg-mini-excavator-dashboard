"""Extract selected jobsite pictures for the formal tonnage dashboards.

The source presentation remains the build-time evidence source. Only selected
picture shapes are exported; full-slide screenshots are never embedded in the
formal dashboards.
"""

from __future__ import annotations

import argparse
import io
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SELECTIONS = {
    "excavator-1-2t": [(18, 3), (20, 3), (22, 2), (24, 1)],
    "excavator-2-3t": [(36, 2), (37, 3), (37, 2)],
    "excavator-35t": [(50, 2), (52, 1), (53, 2), (54, 3)],
    "excavator-4-5t": [(71, 4), (72, 3), (73, 2), (74, 2)],
    "excavator-5-6t": [(92, 1), (93, 1), (94, 2), (95, 1)],
    "excavator-8-10t": [(110, 2), (111, 2), (112, 1), (113, 1)],
    "excavator-12-16t": [(133, 2), (134, 3), (135, 2), (136, 1)],
    "excavator-19-24t": [(154, 1), (155, 2), (157, 2)],
    "excavator-24-33t": [(172, 2), (173, 2), (174, 2), (175, 2)],
    "excavator-33-40t": [(201, 1), (202, 1), (203, 1), (204, 2)],
    "excavator-40-60t": [(218, 1), (219, 2), (220, 3), (221, 1)],
}


def crop_picture(shape) -> Image.Image:
    image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
    width, height = image.size
    left = round(width * shape.crop_left)
    top = round(height * shape.crop_top)
    right = round(width * (1 - shape.crop_right))
    bottom = round(height * (1 - shape.crop_bottom))
    if right > left and bottom > top:
        image = image.crop((left, top, right, bottom))
    image.thumbnail((1280, 900), Image.Resampling.LANCZOS)
    return image


def picture_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def extract(source: Path, output_dir: Path, manifest_path: Path) -> None:
    presentation = Presentation(source)
    output_dir.mkdir(parents=True, exist_ok=True)
    records = []

    for slug, selections in SELECTIONS.items():
        for sequence, (slide_number, picture_number) in enumerate(selections, start=1):
            pictures = picture_shapes(presentation.slides[slide_number - 1])
            if picture_number < 1 or picture_number > len(pictures):
                raise ValueError(
                    f"Slide {slide_number} has {len(pictures)} pictures; "
                    f"picture {picture_number} was requested"
                )
            image = crop_picture(pictures[picture_number - 1])
            filename = f"{slug}-{sequence:02d}.webp"
            image.save(output_dir / filename, "WEBP", quality=84, method=6)
            records.append(
                {
                    "slug": slug,
                    "sequence": sequence,
                    "source_slide": slide_number,
                    "source_picture": picture_number,
                    "file": f"assets/ppt-insights/{filename}",
                    "pixel_width": image.width,
                    "pixel_height": image.height,
                }
            )

    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(
        json.dumps(
            {
                "source": source.name,
                "rule": "Selected picture shapes only; no full-slide screenshots.",
                "records": records,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    print(f"Exported {len(records)} pictures to {output_dir}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output_dir", type=Path)
    parser.add_argument("manifest", type=Path)
    args = parser.parse_args()
    extract(args.source, args.output_dir, args.manifest)


if __name__ == "__main__":
    main()
