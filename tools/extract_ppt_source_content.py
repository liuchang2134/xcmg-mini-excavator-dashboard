"""Extract source-faithful text, tables and visuals from the excavator deck.

The resulting data drives the formal category and tonnage pages. Text is kept
in its original order, table cells remain verbatim, and visual assets are
either exported directly from picture shapes or cropped from PowerPoint's
native slide rendering for charts.
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
from collections import defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DIR = ROOT / "data" / "source-presentations"
OUTPUT_PATH = ROOT / "data" / "ppt-insights" / "ppt-source-content.json"
ASSET_DIR = ROOT / "assets" / "ppt-source"

OVERVIEW_SLIDES = [
    *range(3, 16),
    *range(233, 245),
    246,
]

SLUG_SLIDE_RANGES = {
    "excavator-1-2t": [(16, 34)],
    "excavator-2-3t": [(35, 47)],
    "excavator-35t": [(48, 68)],
    "excavator-4-5t": [(69, 89)],
    "excavator-5-6t": [(90, 107)],
    "excavator-8-10t": [(108, 125)],
    "excavator-12-14t": [(126, 151)],
    "excavator-14-16t-short-tail": [(126, 151)],
    "excavator-21-24t": [(152, 168)],
    "excavator-24-28t": [(169, 187)],
    "excavator-24-28t-short-tail": [(169, 187)],
    "excavator-28-33t": [(169, 177), (188, 198)],
    "excavator-33-40t": [(199, 215)],
    "excavator-40-60t": [(216, 232)],
}

SKIP_EXACT = {
    "二、大区产业板块洞察",
    "二、细分市场洞察",
    "三、市场洞察资源",
    "挖机产品线",
    "销量分析",
    "适应性分析",
    "参配对比",
    "核心性能",
    "产品定位",
}


def normalize_text(value: str) -> str:
    value = (value or "").replace("\r", "\n").replace("\x0b", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r" *\n *", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def source_pptx() -> Path:
    files = sorted(SOURCE_DIR.glob("*.pptx"), key=lambda item: item.stat().st_size, reverse=True)
    if not files:
        raise FileNotFoundError(f"No presentation found in {SOURCE_DIR}")
    return files[0]


def iter_shapes(shape_collection):
    for shape in shape_collection:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return normalize_text(shape.text)


def all_text_items(slide, slide_width, slide_height):
    items = []
    for shape in iter_shapes(slide.shapes):
        value = shape_text(shape)
        if not value:
            continue
        items.append(
            {
                "text": value,
                "x": round(shape.left / slide_width, 6),
                "y": round(shape.top / slide_height, 6),
                "w": round(shape.width / slide_width, 6),
                "h": round(shape.height / slide_height, 6),
            }
        )
    return sorted(items, key=lambda item: (item["y"], item["x"]))


def choose_title(items, slide_number):
    preferred_tokens = (
        "宏观环境分析",
        "行业分析",
        "占有率分析",
        "销售结构分析",
        "核心规格产品",
        "适应性分析",
        "参数及配置对比",
        "产品定位",
        "可售型谱",
        "竞争力提升举措",
        "未来三年",
        "市场洞察资源",
    )
    candidates = []
    for item in items:
        value = item["text"]
        if value in SKIP_EXACT or len(value) > 180:
            continue
        score = 0
        if any(token in value for token in preferred_tokens):
            score += 100
        if re.match(r"^\d+(?:\.\d+)?[、. ]", value):
            score += 40
        if item["y"] < 0.2:
            score += 20
        score -= len(value) / 1000
        candidates.append((score, item["y"], item["x"], value))
    if not candidates:
        return f"第 {slide_number} 页"
    candidates.sort(key=lambda item: (-item[0], item[1], item[2]))
    return candidates[0][3]


def clean_title(value):
    value = normalize_text(value)
    value = re.sub(r"^2、核心产品线[—-]+(?:履带式)?挖掘机\s*", "", value)
    value = re.sub(r"^\d+(?:\.\d+)?\s*", "", value)
    return value.strip(" -—|") or value


def is_body_text(value, title):
    stripped = normalize_text(value)
    if not stripped or stripped == title or stripped in SKIP_EXACT:
        return False
    if re.fullmatch(r"\d{1,3}", stripped):
        return False
    if stripped.startswith(("二、大区产业板块洞察", "二、细分市场洞察")):
        return False
    if stripped == "挖机产品线":
        return False
    return True


def split_source_paragraphs(value):
    value = normalize_text(value)
    if not value:
        return []
    paragraphs = [item.strip() for item in re.split(r"\n{2,}", value) if item.strip()]
    return paragraphs or [value]


def classify_overview(slide_number, title):
    if 3 <= slide_number <= 8:
        return "environment"
    if 9 <= slide_number <= 11:
        return "industry"
    if 12 <= slide_number <= 14:
        return "competition"
    if slide_number == 15:
        return "class_structure"
    if slide_number in {233, 234, 242, 243}:
        return "portfolio"
    if 235 <= slide_number <= 241:
        return "roadmap"
    if slide_number == 244:
        return "sales_plan"
    if slide_number == 246:
        return "intelligence"
    return "overview"


def classify_tonnage(title):
    if any(token in title for token in ("产品销量分析", "细分市场销量分析")):
        return "market"
    if "适应性分析" in title:
        return "applications"
    if "参数及配置对比" in title:
        return "comparison"
    if any(token in title for token in ("产品定位", "占有率目标")):
        return "positioning"
    return "analysis"


def slugs_for_slide(slide_number):
    result = []
    for slug, ranges in SLUG_SLIDE_RANGES.items():
        if any(start <= slide_number <= end for start, end in ranges):
            result.append(slug)
    return result


def picture_image(shape):
    image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
    width, height = image.size
    left = round(width * shape.crop_left)
    top = round(height * shape.crop_top)
    right = round(width * (1 - shape.crop_right))
    bottom = round(height * (1 - shape.crop_bottom))
    if right > left and bottom > top:
        image = image.crop((left, top, right, bottom))
    image.thumbnail((1600, 1100), Image.Resampling.LANCZOS)
    return image


def crop_rendered_shape(rendered_slide, shape, slide_width, slide_height):
    image = rendered_slide
    left = max(0, round(shape.left / slide_width * image.width) - 8)
    top = max(0, round(shape.top / slide_height * image.height) - 8)
    right = min(image.width, round((shape.left + shape.width) / slide_width * image.width) + 8)
    bottom = min(image.height, round((shape.top + shape.height) / slide_height * image.height) + 8)
    if right <= left or bottom <= top:
        return None
    result = image.crop((left, top, right, bottom)).convert("RGB")
    result.thumbnail((1600, 1100), Image.Resampling.LANCZOS)
    return result


def image_digest(image):
    buffer = io.BytesIO()
    image.save(buffer, "PNG", optimize=True)
    return hashlib.sha1(buffer.getvalue()).hexdigest()


def chart_series_text(series):
    value = str(series.name or "").strip()
    if value:
        return value
    element = series._element
    tx = element.find(qn("c:tx"))
    if tx is None:
        return ""
    direct = tx.find(qn("c:v"))
    if direct is not None and direct.text:
        return normalize_text(direct.text)
    for node in tx.iter(qn("c:v")):
        if node.text:
            return normalize_text(node.text)
    return ""


def chart_numeric_points(series, child_name):
    container = series._element.find(qn(f"c:{child_name}"))
    if container is None:
        return []
    indexed = {}
    for point in container.iter(qn("c:pt")):
        value = point.find(qn("c:v"))
        if value is None or value.text in (None, ""):
            continue
        try:
            indexed[int(point.get("idx", len(indexed)))] = float(value.text)
        except (TypeError, ValueError):
            continue
    if indexed:
        return [indexed.get(index) for index in range(max(indexed) + 1)]
    values = []
    for value in container.iter(qn("c:v")):
        try:
            values.append(float(value.text))
        except (TypeError, ValueError):
            continue
    return values


def chart_axis_titles(chart):
    titles = []
    for axis_name in ("catAx", "valAx", "dateAx"):
        for axis in chart._element.iter(qn(f"c:{axis_name}")):
            title = axis.find(qn("c:title"))
            if title is None:
                continue
            text = normalize_text("".join(node.text or "" for node in title.iter(qn("a:t"))))
            if text and text not in titles:
                titles.append(text)
    return titles


def chart_payload(shape):
    """Keep native chart values so supported charts can be rebuilt responsively."""
    chart = shape.chart
    categories = []
    try:
        categories = [str(value) for value in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError, ValueError):
        categories = []
    series = []
    for item in chart.series:
        try:
            values = [
                None if value is None else float(value)
                for value in item.values
            ]
        except (AttributeError, TypeError, ValueError):
            values = []
        record = {
            "name": chart_series_text(item),
            "values": values,
        }
        x_values = chart_numeric_points(item, "xVal")
        bubble_sizes = chart_numeric_points(item, "bubbleSize")
        if x_values:
            record["x_values"] = x_values
        if bubble_sizes:
            record["bubble_sizes"] = bubble_sizes
        series.append(record)
    title = ""
    if chart.has_title:
        title = normalize_text(chart.chart_title.text_frame.text)
    return {
        "type": str(chart.chart_type),
        "title": title,
        "categories": categories,
        "series": series,
        "axis_titles": chart_axis_titles(chart),
    }


def enrich_existing_chart_data():
    payload = json.loads(OUTPUT_PATH.read_text(encoding="utf-8"))
    deck = Presentation(source_pptx())
    records = {int(record["slide"]): record for record in payload.get("slides", [])}
    updated = 0
    for slide_number, slide in enumerate(deck.slides, start=1):
        record = records.get(slide_number)
        if record is None:
            continue
        chart_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_chart", False)
            and (shape.width / deck.slide_width) * (shape.height / deck.slide_height) >= 0.02
        ]
        chart_visuals = [
            visual
            for visual in record.get("visuals", [])
            if visual.get("kind") == "chart"
        ]
        if len(chart_shapes) != len(chart_visuals):
            raise RuntimeError(
                f"Slide {slide_number}: {len(chart_shapes)} chart shapes do not match "
                f"{len(chart_visuals)} chart visuals"
            )
        for visual, shape in zip(chart_visuals, chart_shapes):
            visual["chart_data"] = chart_payload(shape)
            updated += 1
    OUTPUT_PATH.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return updated


def export_visuals(slide, slide_number, rendered_path, slide_width, slide_height, seen):
    rendered = Image.open(rendered_path).convert("RGB") if rendered_path.exists() else None
    visuals = []
    counters = defaultdict(int)
    for shape in slide.shapes:
        area = (shape.width / slide_width) * (shape.height / slide_height)
        image = None
        kind = ""
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            and area >= 0.015
            and shape.width / slide_width >= 0.10
            and shape.height / slide_height >= 0.08
        ):
            image = picture_image(shape)
            kind = "picture"
        elif getattr(shape, "has_chart", False) and area >= 0.02 and rendered is not None:
            image = crop_rendered_shape(rendered, shape, slide_width, slide_height)
            kind = "chart"
        if image is None or image.width < 120 or image.height < 80:
            continue
        digest = image_digest(image)
        counters[kind] += 1
        filename = f"s{slide_number:03d}-{kind}-{counters[kind]:02d}.webp"
        if digest in seen:
            filename = seen[digest]
        else:
            image.save(ASSET_DIR / filename, "WEBP", quality=88, method=6)
            seen[digest] = filename
        visuals.append(
            {
                "kind": kind,
                "file": f"assets/ppt-source/{filename}",
                "width": image.width,
                "height": image.height,
                "x": round(shape.left / slide_width, 6),
                "y": round(shape.top / slide_height, 6),
                "w": round(shape.width / slide_width, 6),
                "h": round(shape.height / slide_height, 6),
                **({"chart_data": chart_payload(shape)} if kind == "chart" else {}),
            }
        )
    if rendered is not None:
        rendered.close()
    return visuals


def load_table_index():
    path = ROOT / "data" / "ppt-insights" / "ppt-business-tables.json"
    payload = json.loads(path.read_text(encoding="utf-8"))
    by_slide = defaultdict(list)
    for record in payload.get("records", []):
        by_slide[int(record["slide"])].append(record["id"])
    return by_slide


def extract(render_dir: Path):
    pptx_path = source_pptx()
    deck = Presentation(pptx_path)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    table_index = load_table_index()
    slides = []
    by_slug = defaultdict(list)
    overview = []
    seen = {}

    for slide_number, slide in enumerate(deck.slides, start=1):
        if slide_number in {1, 2, 245}:
            continue
        slug_values = slugs_for_slide(slide_number)
        is_overview = slide_number in OVERVIEW_SLIDES
        if not slug_values and not is_overview:
            continue

        items = all_text_items(slide, deck.slide_width, deck.slide_height)
        source_title = choose_title(items, slide_number)
        body_blocks = []
        source_notes = []
        for item in items:
            if not is_body_text(item["text"], source_title):
                continue
            for paragraph in split_source_paragraphs(item["text"]):
                target = source_notes if paragraph.startswith(("数据来源", "备注：", "备注:")) else body_blocks
                target.append(paragraph)

        rendered_path = render_dir / f"幻灯片{slide_number}.PNG"
        if not rendered_path.exists():
            rendered_path = render_dir / f"Slide{slide_number}.PNG"
        visuals = export_visuals(
            slide,
            slide_number,
            rendered_path,
            deck.slide_width,
            deck.slide_height,
            seen,
        )
        section = (
            classify_overview(slide_number, source_title)
            if is_overview
            else classify_tonnage(source_title)
        )
        record = {
            "id": f"slide-{slide_number:03d}",
            "slide": slide_number,
            "title": {
                "zh": clean_title(source_title),
                "en": "",
            },
            "source_title_zh": source_title,
            "section": section,
            "body": [{"zh": value, "en": ""} for value in body_blocks],
            "notes": [{"zh": value, "en": ""} for value in source_notes],
            "visuals": visuals,
            "table_ids": table_index.get(slide_number, []),
            "slugs": slug_values,
            "overview": is_overview,
        }
        slides.append(record)
        if is_overview:
            overview.append(record["id"])
        for slug in slug_values:
            by_slug[slug].append(record["id"])

    payload = {
        "meta": {
            "source_file": pptx_path.name,
            "slide_count": len(deck.slides),
            "included_slide_count": len(slides),
            "classification": "XCMG ARC INTERNAL",
            "content_rule": "Original slide text order, original table cells and source visual objects.",
            "excluded": {
                "cover_and_section_slides": [1, 2],
                "personnel_slide": [245],
            },
        },
        "summary": {
            "overview_slides": len(overview),
            "tonnage_slide_links": sum(len(items) for items in by_slug.values()),
            "visual_count": sum(len(record["visuals"]) for record in slides),
            "table_count": sum(len(record["table_ids"]) for record in slides),
            "text_block_count": sum(len(record["body"]) for record in slides),
        },
        "overview": overview,
        "by_slug": dict(by_slug),
        "slides": slides,
    }
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    OUTPUT_PATH.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return payload


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--render-dir",
        type=Path,
        help="PowerPoint-exported PNG slide directory.",
    )
    parser.add_argument(
        "--chart-data-only",
        action="store_true",
        help="Refresh native chart series, axes and bubble values without re-exporting visuals.",
    )
    args = parser.parse_args()
    if args.chart_data_only:
        print(json.dumps({"updated_chart_records": enrich_existing_chart_data()}))
        return
    if args.render_dir is None:
        parser.error("--render-dir is required unless --chart-data-only is used")
    payload = extract(args.render_dir)
    print(json.dumps(payload["summary"], ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
