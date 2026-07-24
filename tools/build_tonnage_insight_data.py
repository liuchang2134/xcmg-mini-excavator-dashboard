"""Build curated tonnage insight data from the internal excavator deck.

Market series are read directly from the embedded PowerPoint charts. Narrative
content is curated by tonnage so one class never inherits another class's
product conclusions. The output is build-time data for the formal dashboards;
it does not change the existing Excel-based scoring model.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def bi(zh: str, en: str) -> dict[str, str]:
    return {"zh": zh, "en": en}


def app(title_zh, title_en, body_zh, body_en, req_zh, req_en, image, slide):
    return {
        "title": bi(title_zh, title_en),
        "body": bi(body_zh, body_en),
        "requirements": bi(req_zh, req_en),
        "image": image,
        "source_slide": slide,
    }


def row(dimension_zh, dimension_en, basis_zh, basis_en, gap_zh, gap_en, action_zh, action_en, slides):
    return {
        "dimension": bi(dimension_zh, dimension_en),
        "basis": bi(basis_zh, basis_en),
        "gap": bi(gap_zh, gap_en),
        "action": bi(action_zh, action_en),
        "source_slides": slides,
    }


INSIGHTS = {
    "excavator-1-2t": {
        "label": bi("1-2 吨级", "1-2 t class"),
        "model": "XCMG XE19U",
        "source_slides": [16, 34],
        "market_slide": 16,
        "concentration": bi("四个领先品牌合计约 65%", "Four leading brands hold about 65%"),
        "benchmark": bi("久保田 U17 / KX018", "Kubota U17 / KX018"),
        "role": bi("狭窄空间、高频转场和轻型运输是本吨级的核心价值，客户同时比较驾驶棚/驾驶室、属具覆盖和非专业机手易用性。", "Confined access, frequent relocation and light transport define this class. Buyers also compare canopy/cab availability, attachment coverage and ease of use for occasional operators."),
        "transport": bi("历史运输组合为轻型皮卡加 6K 拖车；XE18U 自重 1,795 kg。实际交付应按快换、拇指钳、铲斗和驾驶室等完整配置复核组合总重。", "The historical reference is a light-duty pickup with a 6K trailer and a 1,795 kg XE18U base machine. Actual delivery packages must be checked with coupler, thumb, buckets and cab equipment included."),
        "customers": bi("园林承包商、租赁商、市政维护单位、农场及小型施工队", "Landscape contractors, rental fleets, municipalities, farms and small contractors"),
        "jobs": bi("景观整形、住宅改造、管线沟槽、市政维护、农业清理", "Landscaping, residential renovation, utility trenching, municipal maintenance and farm cleanup"),
        "purchase": bi("运输便利、窄机身与短尾、可伸缩底盘、动臂偏摆、驾驶室选择和属具流量", "Transportability, narrow packaging, short tail, variable undercarriage, boom swing, cab choice and auxiliary flow"),
        "applications": [
            app("园林绿化", "Landscaping", "覆盖场地清理、种植坑开挖、地形修整和排水沟施工，要求机器能穿过窄门并减少对既有景观的扰动。", "Covers site cleanup, planting holes, shaping and drainage trenches, with emphasis on passing narrow gates and limiting disturbance to finished areas.", "窄机身、可伸缩底盘、精细微动和低接地扰动", "Narrow width, variable undercarriage, fine modulation and low ground disturbance", "assets/ppt-insights/excavator-1-2t-01.webp", 18),
            app("住宅与商业改造", "Residential and commercial renovation", "用于后院、建筑周边和室内外改造中的开挖、拆除、装车与回填，作业空间通常受围墙和既有建筑约束。", "Supports excavation, light demolition, loading and backfill around homes and commercial buildings where fences and existing structures constrain movement.", "短尾、动臂偏摆、拇指钳和可控回转", "Short tail, boom swing, thumb and controlled swing", "assets/ppt-insights/excavator-1-2t-02.webp", 20),
            app("市政维护", "Municipal maintenance", "面向人行道、路缘、绿化带和浅层管线维护，频繁转场要求启停平顺、装卸快速且安全警示完整。", "Targets sidewalks, curbs, verges and shallow utilities; frequent moves require smooth travel, quick loading and complete warning equipment.", "行走线性、运输效率、报警和照明", "Travel modulation, transport efficiency, alarms and lighting", "assets/ppt-insights/excavator-1-2t-03.webp", 22),
            app("农业与土地维护", "Agriculture and land maintenance", "用于小型沟渠、围栏、灌溉和土地清理，非专业机手更加关注操作直观、稳定性和附件切换。", "Used for small ditches, fencing, irrigation and land cleanup, where occasional operators value intuitive control, stability and attachment changes.", "易操作性、稳定性、快换和多属具", "Ease of use, stability, coupler and multi-attachment capability", "assets/ppt-insights/excavator-1-2t-04.webp", 24),
        ],
        "engineering": [
            row("整机与运输", "Machine and transport", "短尾布置满足受限空间，基础运输质量适合轻型组合。", "Short-tail packaging supports confined work and the base mass fits a light transport concept.", "完整附件包后的组合重量边界尚需统一管理。", "The transport margin of the fully equipped package is not yet governed as one envelope.", "建立基础机、快换、拇指钳、铲斗和驾驶室的配置质量清单。", "Create a configuration mass sheet covering base machine, coupler, thumb, buckets and cab.", [17, 25]),
            row("作业与液压", "Performance and hydraulics", "可选长斗杆、AUX1 流量、斗杆力和系统压力具备纸面基础。", "Optional long stick, AUX1 flow, arm force and system pressure provide a sound paper baseline.", "动臂偏摆范围、行驶速度和复合动作仍弱于标杆。", "Boom-swing range, travel speed and combined-function performance trail the benchmark.", "联合优化偏摆包络、行走标定和复合动作流量分配。", "Jointly optimize boom-swing envelope, travel calibration and combined-flow allocation.", [26, 27, 33]),
            row("操控与安全", "Controls and safety", "安全标识及多数间接安全配置较完整。", "Safety labels and most indirect safety provisions are broadly complete.", "脚踏阀死区、行走加速线性和回转制动需要改进。", "Pedal deadband, travel acceleration linearity and swing braking need improvement.", "开展仪器化操控测试，锁定脚踏阀、回转制动和软件标定目标。", "Run instrumented control tests and lock targets for pedal valves, swing braking and software calibration.", [28, 30]),
            row("驾驶与耐久", "Operator environment and durability", "低油耗和基础维修性具备优势。", "Low fuel use and basic serviceability are positive foundations.", "驾驶室覆盖、涂层防锈和橡胶件耐久弱于标杆。", "Cab availability, coating corrosion resistance and rubber durability trail the benchmark.", "补齐驾驶室方案并实施涂层、紧固件和橡胶件耐久验证。", "Complete the cab offer and validate coatings, fasteners and rubber durability.", [29, 31, 32]),
        ],
    },
    "excavator-2-3t": {
        "label": bi("2-3 吨级", "2-3 t class"), "model": "XCMG XE27U", "source_slides": [35, 47], "market_slide": 35,
        "concentration": bi("四个领先品牌合计超过 80%", "Four leading brands exceed 80% combined"),
        "benchmark": bi("久保田 U27-4 / 山猫同级产品", "Kubota U27-4 / Bobcat class benchmark"),
        "role": bi("本吨级延续微挖的园林、市政和住宅需求，但客户开始更重视驾驶室、挖掘力、AUX 流量和连续作业舒适性。", "This class retains landscaping, municipal and residential demand while placing greater weight on cab availability, digging force, auxiliary flow and all-day comfort."),
        "transport": bi("运输评价应按驾驶室、快换、拇指钳和常用铲斗后的交付质量进行，不仅比较基础机重量；同一拖车组合还需核验轴荷与固定方案。", "Transport should be evaluated at delivered mass with cab, coupler, thumb and common buckets, not base-machine mass alone; axle loading and securement also require verification."),
        "customers": bi("园林与景观承包商、租赁商、市政单位、住宅改造和小型公用工程客户", "Landscape contractors, rental fleets, municipalities, residential renovators and light utility contractors"),
        "jobs": bi("清表、平地、沟槽、破碎、钻孔、抓取、起吊和回填", "Clearing, grading, trenching, breaking, augering, grabbing, lifting and backfill"),
        "purchase": bi("驾驶室与空调、AUX2、多属具适配、行走效率、铲斗力和维修可达性", "Cab and HVAC, AUX2, attachment coverage, travel productivity, bucket force and service access"),
        "applications": [
            app("园林与景观", "Landscaping", "2-3 吨机型可承担更连续的挖沟、整形和材料搬运，同时仍需保持后院和城市绿地的通过性。", "The 2-3 t class supports more continuous trenching, shaping and material handling while retaining access to yards and urban green space.", "短尾、驾驶室选择、微动和多属具", "Short tail, cab choice, fine control and multi-attachment capability", "assets/ppt-insights/excavator-2-3t-01.webp", 36),
            app("市政与管线", "Municipal and utility work", "用于浅层管线、路缘和道路维护，作业链同时包含挖掘、装车、回填和短距离转场。", "Used for shallow utilities, curbs and road maintenance with excavation, loading, backfill and short relocations in one cycle.", "行走速度、AUX 流量、平地协调和警示配置", "Travel speed, auxiliary flow, grading coordination and warning equipment", "assets/ppt-insights/excavator-2-3t-02.webp", 37),
            app("轻拆与属具作业", "Light demolition and attachments", "破碎锤、拇指钳和螺旋钻等属具扩展了应用，但对 AUX1 流量、AUX2、快换和冷却余量提出更高要求。", "Breaker, thumb and auger work expands the application envelope while increasing demand on AUX1 flow, AUX2, coupler performance and cooling reserve.", "AUX1/AUX2、快换、流量设定和驾驶防护", "AUX1/AUX2, coupler, flow settings and operator protection", "assets/ppt-insights/excavator-2-3t-03.webp", 37),
        ],
        "engineering": [
            row("整机参数", "Machine parameters", "基础尺寸和工作范围与久保田总体接近。", "Base dimensions and working range are broadly close to Kubota.", "整机高度与交付质量需按运输组合继续控制。", "Overall height and delivered mass require continued transport-envelope control.", "建立交付状态运输边界并按驾驶室/驾驶棚分别确认。", "Define delivered transport envelopes separately for cab and canopy versions.", [38]),
            row("液压与作业", "Hydraulics and performance", "标配 AUX2，平地复合协调性具备优势。", "Standard AUX2 and grading coordination are strengths.", "行驶速度、AUX1 流量和铲斗挖掘力低于标杆。", "Travel speed, AUX1 flow and bucket digging force trail the benchmark.", "优化泵阀匹配、行走标定和铲斗机构参数。", "Optimize pump/valve matching, travel calibration and bucket linkage targets.", [39, 40, 46]),
            row("操控与驾驶", "Controls and operator environment", "7 英寸显示器和信息呈现具备基础。", "The 7-inch display and information presentation provide a good base.", "双速切换、脚踏死区、空调和座椅减震存在差距。", "Two-speed switching, pedal deadband, HVAC and seat suspension remain gaps.", "把双速控制移至高频操作区，补齐冷暖空调并优化座椅与行走微动。", "Move two-speed control into the primary control zone, add heating/cooling and improve seat and travel modulation.", [43, 44]),
            row("维修与品质", "Serviceability and quality", "燃油经济性优于历史标杆评价。", "Fuel economy was rated favorably against the historical benchmark.", "主阀可达性、涂装防锈和售后资料质量不足。", "Main-valve access, coating corrosion resistance and service-document quality are insufficient.", "调整维修开口并建立涂装、标识和售后文件闭环。", "Improve service openings and close the loop on coating, identification and service documentation.", [42, 45]),
        ],
    },
    "excavator-35t": {
        "label": bi("3-4 吨级", "3-4 t class"), "model": "XCMG XE35U", "source_slides": [48, 68], "market_slide": 48,
        "concentration": bi("四个领先品牌合计约 74%", "Four leading brands hold about 74% combined"),
        "benchmark": bi("久保田 U35/KX033 与迪尔 35P", "Kubota U35/KX033 and Deere 35P"),
        "role": bi("3-4 吨是北美高频通用规格，覆盖承包商、租赁商、市政和景观客户；竞争力取决于带附件运输、多工具作业、短循环效率和全生命周期体验。", "The 3-4 t class is a high-use general-purpose size for contractors, rental fleets, municipalities and landscapers. Competitiveness depends on equipped transport, multi-tool work, short-cycle productivity and lifecycle experience."),
        "transport": bi("历史记录中的 XE35U 基础质量为 4,200 kg；加副配重、快换、拇指钳和标准斗后的典型组合约 4,579 kg，已逼近 14K 拖车常见有效载荷边界。", "The historical XE35U base mass is 4,200 kg. A representative package with counterweight, coupler, thumb and bucket is about 4,579 kg, close to the common payload boundary of a 14K trailer."),
        "customers": bi("中小承包商、租赁商、市政及公用事业单位、景观与农业客户", "Small and mid-size contractors, rental fleets, municipalities, utilities, landscapers and agricultural users"),
        "jobs": bi("沟槽、地基、景观、轻拆、回填、吊装和多属具作业", "Trenching, foundations, landscaping, light demolition, backfill, lifting and multi-attachment work"),
        "purchase": bi("带附件运输质量、属具覆盖、循环效率、可靠性、操控舒适性和残值", "Equipped transport mass, attachment coverage, cycle productivity, reliability, operator comfort and residual value"),
        "applications": [
            app("皮卡与 14K 拖车转运", "Pickup and 14K trailer transport", "经销商交付和客户转场通常要求一套运输组合携带主机及常用附件，因此采购判断必须看交付状态总重。", "Dealer delivery and customer relocation often require one transport package to carry the machine and common attachments, so purchase decisions must use delivered mass.", "整机减重、附件质量清单、装载固定和合规复核", "Machine mass reduction, attachment mass sheet, securement and compliance review", "assets/ppt-insights/excavator-35t-01.webp", 50),
            app("房建与轻拆", "Building work and light demolition", "小型拆除和改造需要铲斗、拇指钳、破碎锤与封闭驾驶环境，并要求在受限空间内完成拆分和装车。", "Light demolition and renovation require buckets, thumb, breaker and an enclosed operator environment, with sorting and loading in constrained space.", "AUX 流量、驾驶防护、精细操控和作业范围", "Auxiliary flow, operator protection, fine control and working range", "assets/ppt-insights/excavator-35t-02.webp", 52),
            app("地基与管线", "Foundation and utility excavation", "现场测量、开挖、清底、管线安装和回填形成完整任务链，短尾和动臂偏摆直接影响贴边施工。", "Surveying, excavation, cleanup, utility installation and backfill form a complete workflow; short tail and boom swing directly affect close-to-structure work.", "挖掘深度、偏摆包络、回转力矩和推土铲功能", "Digging depth, boom-swing envelope, swing torque and blade functions", "assets/ppt-insights/excavator-35t-03.webp", 53),
            app("景观与排水", "Landscaping and drainage", "从清表、整形、沟槽到回填和绿化收尾，需要一台机器频繁切换挖掘、搬运、平地和吊装动作。", "From clearing and shaping through trenching, backfill and final landscaping, one machine must switch frequently among digging, handling, grading and lifting.", "多属具、坡度/深度控制、微动和回填效率", "Multi-attachment capability, grade/depth control, fine modulation and backfill productivity", "assets/ppt-insights/excavator-35t-04.webp", 54),
        ],
        "engineering": [
            row("运输与型谱", "Transport and portfolio", "短尾可覆盖受限空间，附件配置较完整。", "The short-tail package supports confined work with broad attachment availability.", "整机偏重，且型谱缺少常规尾重载选择。", "Machine mass is high and the portfolio lacks a conventional-tail heavy-duty choice.", "锁定轻量化交付质量边界，并评估同平台常规尾产品。", "Lock a lightweight delivered-mass envelope and evaluate a conventional-tail platform variant.", [49, 59, 67]),
            row("性能与液压", "Performance and hydraulics", "斗杆力、AUX2 和起吊数据具备纸面基础。", "Arm force, AUX2 and lift data provide paper strengths.", "行驶速度、AUX1 流量、铲斗力和回转回填效率偏弱。", "Travel speed, AUX1 flow, bucket force and swing/backfill productivity are weak.", "联合优化行走、附件液压和回转制动/力矩标定。", "Jointly optimize travel, attachment hydraulics and swing braking/torque calibration.", [60, 61, 67]),
            row("操控与驾驶", "Controls and cab", "安全配置整体接近历史标杆。", "Safety equipment is broadly close to the historical benchmark.", "微动、精准性、平顺性、座椅、空间和显示交互存在差距。", "Inching, precision, smoothness, seat, space and display interaction remain gaps.", "以仪器测试和目标机手盲评锁定阀控、软件及驾驶室目标。", "Use instrumented tests and blind operator evaluations to set valve, software and cab targets.", [62, 64, 65]),
            row("耐久与交付", "Durability and support", "燃油经济性具备潜在优势。", "Fuel economy shows potential advantage.", "防锈、涂层、橡胶件、维修加注和售后资料弱于标杆。", "Corrosion protection, coatings, rubber parts, service filling and after-sales documentation trail the benchmark.", "补齐耐久验证、维修工时和文件准确性闭环。", "Close validation loops for durability, service labor and documentation accuracy.", [63, 66]),
        ],
    },
    "excavator-4-5t": {
        "label": bi("4-5 吨级", "4-5 t class"), "model": "XCMG XE45U", "source_slides": [69, 89], "market_slide": 69,
        "concentration": bi("四个领先品牌合计超过 85%", "Four leading brands exceed 85% combined"),
        "benchmark": bi("久保田 U48 / KX040 与卡特同级产品", "Kubota U48 / KX040 and Caterpillar class benchmark"),
        "role": bi("4-5 吨是北美微小挖销量第三的区间，市场主要集中在 4 吨级，购买逻辑是在作业能力和 14K 运输组合之间寻找平衡。", "The 4-5 t range is the third-largest compact class in the source market view, concentrated near 4 t, where buyers balance work capability against a 14K transport package."),
        "transport": bi("XE45U 加副配重、快换、拇指钳和挖沟斗后约 5.16 t；再携带清理斗时超过资料中的 5.2 t 有效载荷上限，运输灵活性成为核心约束。", "XE45U reaches about 5.16 t with counterweight, coupler, thumb and trenching bucket; adding a cleanup bucket exceeds the 5.2 t reference payload, making transport flexibility a primary constraint."),
        "customers": bi("市政与道路承包商、景观与住宅施工企业、租赁商和公用工程客户", "Municipal and road contractors, landscape and residential builders, rental fleets and utility contractors"),
        "jobs": bi("人行道、排水沟、小型拆除、泳池与地基、景观土方", "Sidewalks, drainage ditches, light demolition, pools and foundations, and landscape earthwork"),
        "purchase": bi("带附件运输质量、稳定与起吊、AUX 流量、驾驶便利和回填效率", "Equipped transport mass, stability and lift, auxiliary flow, operator convenience and backfill productivity"),
        "applications": [
            app("人行道与路缘", "Sidewalks and curbs", "狭窄道路和人行区施工需要小回转包络、精准平地和对周边人员的安全隔离。", "Work on narrow roads and pedestrian areas requires compact swing clearance, precise grading and safe separation from nearby people.", "短尾、摄像头、精细操控和快速转场", "Short tail, camera, fine control and fast relocation", "assets/ppt-insights/excavator-4-5t-01.webp", 71),
            app("排水沟施工", "Drainage ditch work", "开挖、清底、管件安装和回填要求工作范围、AUX 配置与推土铲动作协同。", "Excavation, cleanup, pipe installation and backfill require coordinated working range, auxiliary circuits and blade functions.", "挖深、流量、回填效率和坡度控制", "Digging depth, flow, backfill productivity and grade control", "assets/ppt-insights/excavator-4-5t-02.webp", 72),
            app("小型建筑拆除", "Light building demolition", "破碎、拆分、抓取和装车连续发生，对铲斗力、拇指钳、破碎管路和驾驶防护提出综合要求。", "Breaking, sorting, grabbing and loading occur in one cycle, combining requirements for bucket force, thumb, breaker plumbing and operator protection.", "铲斗力、AUX、快换、驾驶防护", "Bucket force, AUX, coupler and operator protection", "assets/ppt-insights/excavator-4-5t-03.webp", 73),
            app("泳池与地基", "Pool and foundation excavation", "受限宅基地内需要深挖、贴边、清底和土方装运，稳定性和最远端操控直接影响质量。", "Confined residential sites require digging, close-wall work, cleanup and spoil handling; stability and far-reach control directly affect quality.", "工作范围、稳定性、微动和运输组合", "Working range, stability, fine control and transport package", "assets/ppt-insights/excavator-4-5t-04.webp", 74),
        ],
        "engineering": [
            row("运输与稳定", "Transport and stability", "整机稳定性和起吊表现优于历史标杆判断。", "Machine stability and lifting were rated favorably in the historical comparison.", "带完整附件运输余量不足，资料建议减重 100-150 kg。", "The fully equipped transport margin is limited; the source recommends a 100-150 kg mass reduction.", "以 14K 运输包为边界开展轻量化，并保持起吊稳定性。", "Use the 14K transport package as the mass target while preserving lifting stability.", [70, 80]),
            row("性能与属具", "Performance and attachments", "斗杆力、系统流量和机具流量具备优势。", "Arm force, system flow and attachment flow are strengths.", "铲斗挖掘力偏低，回转启停晃动稍大。", "Bucket force is low and swing start/stop oscillation is excessive.", "优化铲斗机构和回转制动标定，并验证复合动作。", "Optimize bucket linkage and swing-braking calibration, then validate combined functions.", [81, 82, 85]),
            row("安全与驾驶", "Safety and cab", "起吊报警、回转制动、后视摄像头和 7 英寸交互具备基础。", "Lift alarm, swing brake, rear camera and 7-inch interface provide a sound base.", "驾驶操作仍需在高频任务中验证可及性和视野。", "Control reach and visibility still require validation in frequent tasks.", "按目标机手和典型任务开展人机诊断。", "Run ergonomic diagnostics with target operators and representative tasks.", [83, 86]),
            row("耐久与品质", "Durability and quality", "覆盖件开启与基础维修性接近标杆。", "Panel access and basic serviceability are close to the benchmark.", "油漆、防锈和橡胶件耐久仍需加强。", "Paint, corrosion protection and rubber durability still need improvement.", "补充高海拔、防腐和长期耐久试验。", "Add high-altitude, corrosion and long-duration validation.", [84, 87, 88]),
        ],
    },
    "excavator-5-6t": {
        "label": bi("5-6 吨级", "5-6 t class"), "model": "XCMG XE55U", "source_slides": [90, 107], "market_slide": 90,
        "concentration": bi("四个领先品牌合计超过 70%", "Four leading brands exceed 70% combined"),
        "benchmark": bi("久保田 KX057/U55 与卡特 305 CR", "Kubota KX057/U55 and Caterpillar 305 CR"),
        "role": bi("5-6 吨是北美微小挖销量第二的区间；常规尾销量领先短尾，说明客户对起吊、稳定和重载能力的权重上升。", "The 5-6 t range is the second-largest compact class in the source view. Conventional-tail volume leads short-tail volume, indicating greater emphasis on lift, stability and heavy-duty capability."),
        "transport": bi("XE55U 带副配重、快换、拇指钳和两只常用铲斗后的典型组合约 6,255 kg，处于 16K 拖车资料参考载荷范围内。", "A representative XE55U package with counterweight, coupler, thumb and two common buckets is about 6,255 kg, within the referenced 16K trailer payload range."),
        "customers": bi("市政与公用工程承包商、景观施工企业、租赁商和中小土方客户", "Municipal and utility contractors, landscapers, rental fleets and small earthmoving contractors"),
        "jobs": bi("沟槽与排水、基础、材料搬运、重载回填、景观和多属具作业", "Trenching and drainage, foundations, material handling, heavy backfill, landscaping and multi-attachment work"),
        "purchase": bi("稳定与起吊、发动机储备、铲斗力、属具流量、推土铲操控和安全摄像头", "Stability and lift, engine reserve, bucket force, attachment flow, blade controls and safety camera"),
        "applications": [
            app("市政沟槽", "Municipal trenching", "持续挖沟、管件搬运和回填要求液压温控、行走效率和推土铲操控连续稳定。", "Continuous trenching, pipe handling and backfill require stable hydraulic temperature, travel productivity and blade control.", "长时间液压性能、行走、吊装和回填", "Sustained hydraulics, travel, lifting and backfill", "assets/ppt-insights/excavator-5-6t-01.webp", 92),
            app("房建与基础", "Building and foundation work", "本吨级承担更大土方量和更重物料，常规尾与副配重能够改善最远端稳定和起吊。", "This class handles larger spoil volumes and heavier materials; conventional-tail packaging and counterweight improve far-reach stability and lift.", "稳定性、工作范围、铲斗力和起吊表", "Stability, working range, bucket force and lift chart", "assets/ppt-insights/excavator-5-6t-02.webp", 93),
            app("景观与整形", "Landscaping and grading", "在景观工程中同时完成清表、挖沟、平地和材料搬运，操作手对微动、双速和偏摆控制非常敏感。", "Landscape projects combine clearing, trenching, grading and material handling, making modulation, two-speed and blade-angle controls highly visible to operators.", "微动、双速、推土铲偏摆和低噪声", "Fine control, two-speed, angle blade and low noise", "assets/ppt-insights/excavator-5-6t-03.webp", 94),
            app("多属具作业", "Multi-attachment work", "高流量附件和两套管路扩大破碎、钻孔和抓取能力，但要求发动机与主泵具备持续功率余量。", "High-flow tools and dual circuits expand breaker, auger and grapple use but require sustained engine and pump power reserve.", "AUX1/AUX2、冷却、发动机储备和快换", "AUX1/AUX2, cooling, engine reserve and coupler", "assets/ppt-insights/excavator-5-6t-04.webp", 95),
        ],
        "engineering": [
            row("稳定与起吊", "Stability and lifting", "副配重、较高稳定性和起吊能力是纸面优势。", "Counterweight, stability and lifting capability are paper strengths.", "当前仅一款短尾，重载型谱覆盖弱于久保田。", "A single short-tail model provides less heavy-duty portfolio coverage than Kubota.", "形成副配重标准包并评估重载/常规尾型谱。", "Standardize a counterweight package and evaluate heavy-duty/conventional-tail coverage.", [90, 98]),
            row("动力与作业", "Power and performance", "斗杆力、起吊和附件流量表现较好。", "Arm force, lift and attachment flow perform well.", "发动机功率储备、铲斗力和行走驱动力偏弱。", "Engine reserve, bucket force and travel drive are weak.", "升级动力与主泵匹配，优化铲斗机构和行走牵引。", "Upgrade engine/pump matching and optimize bucket linkage and travel traction.", [99, 100, 106]),
            row("操控与安全", "Controls and safety", "多数安全功能及基础操控接近标杆。", "Most safety functions and base controls are close to the benchmark.", "缺少后视摄像头；偏摆与双速开关布局不便。", "Rear camera is missing and angle-blade/two-speed switch placement is inconvenient.", "标配摄像头并把偏摆、双速集成到推土铲手柄。", "Standardize the camera and integrate angle-blade/two-speed controls into the blade lever.", [101, 103, 104]),
            row("环境与维修", "Environment and service", "燃油经济性和基础维修性具备优势。", "Fuel economy and basic serviceability are positive.", "高海拔动力储备、独立液压油加注和表面耐久不足。", "High-altitude power reserve, independent hydraulic-oil filling and surface durability are insufficient.", "采用增压动力方案并改善加注口、防腐和橡胶件验证。", "Use a turbocharged power solution and improve fill access, corrosion protection and rubber validation.", [102, 105]),
        ],
    },
    "excavator-7-8t": {
        "label": bi("7-8 吨级", "7-8 t class"), "model": "XCMG XE75U", "source_slides": [11, 15], "market_slide": None,
        "concentration": bi("小挖扩展规格", "Extended compact-excavator class"),
        "benchmark": bi("久保田小挖产品体系", "Kubota compact-excavator portfolio"),
        "role": bi("6-8 吨在产品线销量结构中属于小挖次销规格，价值不在单纯放大参数，而在兼顾运输、稳定、多属具和高利用率。", "The 6-8 t range is an extended compact class. Its value is not simply larger specifications but a balance of transport, stability, attachment versatility and high utilization."),
        "transport": bi("本吨级应以交付状态质量、拖车轴荷、机宽机高和常用附件同车运输能力建立边界，并由最新州法规和承运方案复核。", "This class requires a delivered-mass envelope covering trailer axle loads, machine width/height and common attachments carried together, with current state regulations and carrier plans verified."),
        "customers": bi("租赁商、市政与公用工程承包商、景观与中小土方客户", "Rental fleets, municipal and utility contractors, landscapers and small earthmoving contractors"),
        "jobs": bi("公用工程、沟槽、材料搬运、回填、坡地和多属具作业", "Utility work, trenching, material handling, backfill, slopes and multi-attachment work"),
        "purchase": bi("利用率、稳定与起吊、属具覆盖、运输包、可靠性和操控一致性", "Utilization, stability and lift, attachment coverage, transport package, reliability and consistent controls"),
        "applications": [
            app("市政与管线", "Municipal and utility work", "典型任务包括沟槽、管件吊装、回填和路面恢复，一台机器需要覆盖完整任务链。", "Typical work includes trenching, pipe handling, backfill and surface restoration, requiring one machine to cover the full task chain.", "工作范围、起吊、附件管路和回填效率", "Working range, lift, attachment circuits and backfill productivity", "assets/conditions/trench-foundation.jpg", 15),
            app("材料搬运与装车", "Material handling and loading", "更高质量级带来稳定性，同时也要求回转、动臂提升和装车循环保持协调。", "The higher mass class adds stability while requiring coordinated swing, boom lift and loading cycles.", "稳定性、回转力矩、提升速度和视野", "Stability, swing torque, boom speed and visibility", "assets/conditions/loading-cycle.jpg", 15),
            app("坡地与软土", "Slopes and soft ground", "租赁和市政客户需要机器在坡地、松软地面和不同机手之间保持可预测的行走与制动响应。", "Rental and municipal users need predictable travel and braking on slopes, soft ground and across different operators.", "牵引、制动、接地与操控一致性", "Traction, braking, ground contact and control consistency", "assets/conditions/slope-soft-ground.jpg", 15),
        ],
        "engineering": [
            row("市场角色", "Market role", "该吨级补齐小挖从轻型施工到高利用率工程任务的过渡。", "This class bridges light compact work and higher-utilization contracting tasks.", "专项销量与实机场景资料没有单独拆分。", "Dedicated class-level sales and field-use evidence is not separately reported.", "用当前源表和北美客户样本建立本吨级独立市场与工况证据。", "Build class-specific market and jobsite evidence from the current workbook and North American customer samples.", [11, 15]),
            row("运输与配置", "Transport and equipment", "短尾和多属具能够覆盖市政与租赁需求。", "Short-tail packaging and multi-attachment capability support municipal and rental use.", "交付状态运输包和常用附件组合需要统一。", "Delivered transport and common attachment packages need one definition.", "建立按客户类型区分的运输与附件包。", "Define transport and attachment packages by customer type.", [15]),
            row("操控与效率", "Controls and productivity", "本吨级应延续小挖对易用性和多功能的要求。", "This class should retain compact-machine ease of use and versatility.", "跨机手一致性、短循环效率和坡地操控需要实测。", "Cross-operator consistency, short-cycle productivity and slope control require field tests.", "开展循环计时、坡地行走和多机手盲评。", "Run cycle timing, slope travel and multi-operator blind evaluations.", [11, 15]),
            row("可靠性与残值", "Reliability and residual value", "成熟市场客户高度重视稳定性和使用成本。", "Mature-market customers place high value on stability and ownership cost.", "当前缺少本吨级独立的耐久、维修工时和残值闭环。", "A class-specific loop for durability, service labor and residual value is missing.", "建立故障、保养工时、备件和二手机残值台账。", "Establish failure, service-labor, parts and used-value tracking.", [11, 13]),
        ],
    },
    "excavator-8-10t": {
        "label": bi("8-10 吨级", "8-10 t class"), "model": "XCMG XE80U", "source_slides": [108, 125], "market_slide": 108,
        "concentration": bi("四个领先品牌合计超过 70%", "Four leading brands exceed 70% combined"),
        "benchmark": bi("久保田 KX080 与卡特 308/309 CR", "Kubota KX080 and Caterpillar 308/309 CR"),
        "role": bi("8-10 吨是微小挖销量第四的区间，客户已经接受 CDL 和重型运输组合，因此更关注生产率、稳定、起吊与附件能力。", "The 8-10 t range is the fourth-largest compact segment in the source view. With CDL and heavy transport already accepted, buyers focus more on productivity, stability, lifting and attachment capability."),
        "transport": bi("XE80U 带快换、拇指钳和常用铲斗后的典型组合约 10,210 kg，资料采用重型皮卡加 30K 鹅颈拖车，并明确需要 CDL。", "A representative XE80U package with coupler, thumb and common buckets is about 10,210 kg. The source uses a heavy-duty pickup and 30K gooseneck trailer and identifies CDL operation."),
        "customers": bi("专业承包商、租赁商、市政与公用工程客户、景观和轻土方企业", "Professional contractors, rental fleets, municipal/utility users, landscapers and light earthmoving firms"),
        "jobs": bi("深沟、基础、土方、材料搬运、破碎、起吊和道路维护", "Deep trenching, foundations, earthmoving, material handling, breaking, lifting and road maintenance"),
        "purchase": bi("稳定、起吊、挖掘力、液压流量、运输高度、长斗杆和驾驶舒适性", "Stability, lift, digging forces, hydraulic flow, transport height, long stick and operator comfort"),
        "applications": [
            app("道路与市政", "Road and municipal work", "更长作业循环和更重材料要求机器在挖掘、回转、行走和回填之间保持效率。", "Longer cycles and heavier materials require efficient transitions among digging, swing, travel and backfill.", "循环效率、视野、推土铲和报警", "Cycle productivity, visibility, blade functions and alarms", "assets/ppt-insights/excavator-8-10t-01.webp", 110),
            app("基础与深沟", "Foundations and deep trenches", "深沟和基础开挖对长斗杆、稳定、挖深和吊装能力提出更高要求。", "Deep trench and foundation work places greater demand on long-stick availability, stability, depth and lifting capability.", "长斗杆、挖深、稳定和起吊表", "Long stick, digging depth, stability and lift chart", "assets/ppt-insights/excavator-8-10t-02.webp", 111),
            app("材料搬运", "Material handling", "拇指钳、快换和大流量附件扩展了搬运与拆分任务，但牵引和回转响应需要匹配更大质量。", "Thumb, coupler and high-flow tools expand handling and sorting, while travel and swing response must match the higher mass.", "AUX 流量、牵引、回转和稳定", "Auxiliary flow, traction, swing and stability", "assets/ppt-insights/excavator-8-10t-03.webp", 112),
            app("轻拆与多属具", "Light demolition and multi-tool work", "更高发动机和液压能力支持破碎和切割，驾驶防护、冷却和附件设定决定持续作业能力。", "Higher engine and hydraulic capability supports breaking and cutting; cab protection, cooling and tool settings determine sustained output.", "防护、冷却、AUX 设定和快换", "Protection, cooling, AUX settings and coupler", "assets/ppt-insights/excavator-8-10t-04.webp", 113),
        ],
        "engineering": [
            row("整机与运输", "Machine and transport", "动力和液压基础较强。", "Power and hydraulic foundations are strong.", "整机偏重、运输高度偏高，长斗杆状态下稳定性不足。", "Machine mass and transport height are high, with insufficient stability in long-stick configuration.", "轻量化并重算长斗杆、配重和运输高度组合。", "Reduce mass and recalculate long-stick, counterweight and transport-height combinations.", [109, 116]),
            row("作业性能", "Working performance", "主泵流量、压力和 AUX1 流量具备优势。", "Main-pump flow, pressure and AUX1 flow are strengths.", "铲斗力、牵引和起吊能力弱于主要竞品。", "Bucket force, traction and lift trail key competitors.", "优化机构、行走驱动和稳定性，按任务循环验证。", "Optimize linkage, travel drive and stability, then validate by task cycle.", [117, 118]),
            row("操控与驾驶", "Controls and operator environment", "平地与行走优化后得到较好市场反馈。", "Grading and travel improvements received favorable market feedback.", "双速开关未集成，双动臂版本稳定性仍需提升。", "Two-speed control is not integrated and dual-boom stability still needs improvement.", "优化手柄控制和双动臂配重/底盘匹配。", "Improve lever controls and counterweight/undercarriage matching for dual-boom variants.", [121, 122, 124]),
            row("可靠性与维修", "Reliability and service", "三面开启覆盖件有利于维修。", "Three-side panel access supports serviceability.", "表面耐久和工作灯延迟等细节配置存在差距。", "Surface durability and details such as delayed work lights remain gaps.", "补充耐久验证并建立高频客户配置清单。", "Add durability validation and define a high-frequency customer equipment list.", [119, 120, 123]),
        ],
    },
    "excavator-12-14t": {
        "label": bi("12-14 吨级", "12-14 t class"), "model": "XCMG XE135U", "source_slides": [126, 151], "market_slide": 126,
        "concentration": bi("四个领先品牌合计超过 70%", "Four leading brands exceed 70% combined"),
        "benchmark": bi("卡特同级中挖", "Caterpillar class benchmark"),
        "role": bi("12-14 吨市场总量较低且历史销量下滑，产品决策应优先验证真实客户和应用，而不是仅用参数覆盖证明必要性。", "The 12-14 t market is relatively small and historically declined, so product decisions should validate real customers and applications rather than rely on specification coverage alone."),
        "transport": bi("13-16 吨级资料采用阶梯式拖车；宽度目标为 2,590 mm，并按单轴 9,072 kg 参考值配置两轴运输。实际运输需按州法规复核。", "The 13-16 t reference uses a step-deck trailer, a 2,590 mm width target and two axles based on a 9,072 kg per-axle reference. Actual moves require state-level verification."),
        "customers": bi("市政、公用工程、房建、道路施工和专业租赁客户", "Municipal, utility, building, roadwork and professional rental customers"),
        "jobs": bi("沟槽、基础、土方、吊装、破碎和市政综合施工", "Trenching, foundations, earthmoving, lifting, breaking and mixed municipal work"),
        "purchase": bi("运输尺寸、长斗杆、起吊、行走效率、附件液压和智能辅助", "Transport dimensions, long stick, lifting, travel productivity, attachment hydraulics and machine guidance"),
        "applications": [
            app("城市基础施工", "Urban foundation work", "在城市项目中完成基坑、管线和回填，运输宽度、尾部包络和作业范围共同决定适配性。", "Urban projects combine excavation, utilities and backfill, where transport width, tail envelope and working range jointly determine fit.", "运输宽度、工作范围、稳定和回填效率", "Transport width, working range, stability and backfill productivity", "assets/ppt-insights/excavator-12-16t-01.webp", 133),
            app("深沟与管线", "Deep trench and utilities", "较深管线要求长斗杆、挖深、吊装和精确清底，缺少长斗杆会直接缩小任务覆盖。", "Deeper utilities require a long stick, digging depth, lifting and precise cleanup; lack of a long stick directly narrows task coverage.", "长斗杆、挖深、起吊和 2D 引导", "Long stick, digging depth, lift and 2D guidance", "assets/ppt-insights/excavator-12-16t-02.webp", 134),
            app("土方与装车", "Earthmoving and loading", "重载挖掘和装车需要挖掘力、回转、动臂提升和液压流量协同。", "Heavy digging and truck loading require coordinated digging force, swing, boom lift and hydraulic flow.", "循环时间、回转力矩、流量和燃油", "Cycle time, swing torque, flow and fuel", "assets/ppt-insights/excavator-12-16t-03.webp", 135),
            app("破碎与属具", "Breaker and attachment work", "PTO、AUX 和卸油能力决定破碎锤及其他机具的持续输出，驾驶室防护决定安全边界。", "PTO, AUX and case-drain capability determine sustained breaker and tool output, while cab protection defines the safety envelope.", "PTO/AUX、卸油、冷却和防护", "PTO/AUX, case drain, cooling and protection", "assets/ppt-insights/excavator-12-16t-04.webp", 136),
        ],
        "engineering": [
            row("工作范围", "Working range", "发动机功率和重载挖掘基础较强。", "Engine power and heavy-digging foundation are strong.", "标配斗杆偏短，挖掘范围不足。", "The standard stick is short and limits working range.", "加长标配斗杆并重算稳定、起吊和运输包络。", "Lengthen the standard stick and recalculate stability, lift and transport envelope.", [137, 138]),
            row("液压与附件", "Hydraulics and attachments", "PTO、AUX1 合流和较大流量有利于附件作业。", "PTO, AUX1 combining and high flow support attachment work.", "全电控精度和复合动作智能化弱于卡特。", "Electrohydraulic precision and combined-function intelligence trail Caterpillar.", "建立主泵、主阀和控制器联合标定目标。", "Set joint calibration targets for pump, valve and controller.", [138, 139]),
            row("力学与转场", "Forces and relocation", "挖掘力和行走牵引力具备优势。", "Digging force and travel traction are strengths.", "起吊能力和行走速度弱于标杆。", "Lift and travel speed trail the benchmark.", "优化稳定、起吊和行走驱动并完成载荷表验证。", "Optimize stability, lift and travel drive and validate the lift chart.", [140]),
            row("智能与验证", "Intelligence and validation", "现有安全与附件基础可支持扩展。", "Existing safety and attachment foundations support further development.", "缺少 2D/3D、电子围栏和称重等智能功能。", "2D/3D guidance, e-fence and payload functions are missing.", "按客户任务配置智能包，并用作业效率和安全收益验证。", "Configure an intelligent package by customer task and validate productivity and safety benefit.", [139, 150]),
        ],
    },
    "excavator-14-16t-short-tail": {
        "label": bi("14-16 吨级短尾", "14-16 t short-tail class"), "model": "XCMG XE155UCR", "source_slides": [127, 151], "market_slide": 127,
        "concentration": bi("三个领先品牌合计约 64%", "Three leading brands hold about 64% combined"),
        "benchmark": bi("卡特短尾中挖", "Caterpillar short-tail benchmark"),
        "role": bi("14-16 吨是中挖热门短尾区间，面向城市道路、住宅和公共设施建设，核心价值是受限空间中的中挖生产率。", "The 14-16 t range is a popular short-tail mid-size class for urban roads, residential and public infrastructure, combining confined-site packaging with mid-size productivity."),
        "transport": bi("资料采用阶梯式拖车并控制 2,590 mm 宽度；两轴运输能力按交付质量复核，州际限高与许可要求需逐项目确认。", "The reference uses a step-deck trailer with a 2,590 mm width target. Two-axle capacity must be checked at delivered mass and interstate height/permit requirements confirmed by project."),
        "customers": bi("城市道路、市政、公用工程、住宅建设和专业租赁客户", "Urban road, municipal, utility, residential construction and professional rental customers"),
        "jobs": bi("受限空间基础、管线、土方、吊装、破碎和回填", "Confined foundations, utilities, earthmoving, lifting, breaking and backfill"),
        "purchase": bi("短尾包络、工作范围、起吊、流量、行走平顺和智能辅助", "Short-tail envelope, working range, lift, flow, travel smoothness and machine guidance"),
        "applications": [
            app("城市道路", "Urban roadwork", "短尾降低对车道和周边设施的侵入，但视野、警示和回转制动仍决定现场安全。", "Short-tail packaging reduces lane and infrastructure intrusion, while visibility, warnings and swing braking still determine jobsite safety.", "短尾、摄像头、警示和回转控制", "Short tail, cameras, warnings and swing control", "assets/ppt-insights/excavator-12-16t-01.webp", 133),
            app("深沟与基础", "Deep trench and foundations", "在受限空间中同时要求挖深、长斗杆和稳定，单纯缩短尾部不能替代工作范围。", "Confined sites still require depth, long-stick coverage and stability; short tail alone cannot replace working range.", "长斗杆、挖深、稳定和精确控制", "Long stick, depth, stability and precise control", "assets/ppt-insights/excavator-12-16t-02.webp", 134),
            app("吊装与回填", "Lifting and backfill", "管件和预制件吊装后紧接回填与整平，要求起吊表、回转、推土铲和行走复合稳定。", "Pipe and precast lifting is followed by backfill and grading, requiring a verified lift chart and stable swing, blade and travel combinations.", "起吊、回转、推土铲和复合动作", "Lift, swing, blade and combined functions", "assets/ppt-insights/excavator-12-16t-03.webp", 135),
            app("破碎与多属具", "Breaker and multi-tool work", "城市改造对低噪声、驾驶防护、附件管路和快速换装提出更高要求。", "Urban renovation raises requirements for low noise, operator protection, attachment plumbing and rapid tool changes.", "AUX、卸油、防护、冷却和快换", "AUX, case drain, protection, cooling and coupler", "assets/ppt-insights/excavator-12-16t-04.webp", 136),
        ],
        "engineering": [
            row("短尾与范围", "Short tail and range", "短尾布置适合城市受限工地。", "Short-tail packaging fits confined urban jobsites.", "标配斗杆偏短，挖掘范围没有优势。", "The standard stick is short and working range lacks advantage.", "加长斗杆并同步验证尾部、稳定和运输边界。", "Lengthen the stick and jointly validate tail clearance, stability and transport limits.", [141, 150]),
            row("动力与液压", "Power and hydraulics", "发动机、挖掘力、起吊和附件系统基础较好。", "Engine power, digging force, lift and attachment system provide a good base.", "主泵流量和全电控智能化弱于标杆。", "Main-pump flow and electrohydraulic intelligence trail the benchmark.", "升级主泵并建立泵阀控一体化标定。", "Upgrade the main pump and establish integrated pump/valve/controller calibration.", [142, 143, 144]),
            row("行走与复合", "Travel and combined functions", "行走电控与阀芯优化已提供改进基础。", "Travel electronic control and spool optimization provide an improvement base.", "行走与上车复合的平顺性仍有差距。", "Travel and upper-structure combined-function smoothness remains a gap.", "用加速度、冲击和循环时间建立客观标定目标。", "Set objective calibration targets using acceleration, shock and cycle time.", [134, 150]),
            row("智能辅助", "Machine guidance", "现有安全和机具系统可作为集成平台。", "Existing safety and tool systems can support integration.", "2D/3D 智能辅助和高级安全功能不足。", "2D/3D guidance and advanced safety functions are insufficient.", "按城市沟槽、平地和吊装任务配置智能包。", "Configure an intelligent package for urban trenching, grading and lifting.", [143, 150]),
        ],
    },
    "excavator-21-24t": {
        "label": bi("21-24 吨级", "21-24 t class"), "model": "XCMG XE225U", "source_slides": [152, 168], "market_slide": 152,
        "concentration": bi("三个领先品牌合计约 70%", "Three leading brands hold about 70% combined"),
        "benchmark": bi("卡特 320/323 产品体系", "Caterpillar 320/323 portfolio"),
        "role": bi("19-24 吨是中大挖热门区间，覆盖中大型土方、建筑地基和市政工程；竞品通过高、中、经济配置形成客户分层。", "The 19-24 t range is a high-volume mid-size class for earthmoving, foundations and municipal work. Competitors segment customers with premium, core and economy configurations."),
        "transport": bi("XE225U 运输高度 3,030 mm，资料采用阶梯式拖车和三轴承载，并指出机宽超过 2,590 mm 时需要超宽许可。", "XE225U transport height is 3,030 mm. The reference uses a step-deck trailer and three axles and identifies an oversize-width permit when machine width exceeds 2,590 mm."),
        "customers": bi("土方与基础承包商、市政道路企业、专业租赁商和公用工程客户", "Earthmoving and foundation contractors, municipal road firms, professional rental fleets and utilities"),
        "jobs": bi("土方装车、基础、沟槽、破碎、吊装和场地平整", "Truck loading, foundations, trenching, breaking, lifting and site grading"),
        "purchase": bi("循环效率、可靠性、附件系统、起吊、智能辅助、服务与残值", "Cycle productivity, reliability, attachment system, lifting, machine guidance, service and residual value"),
        "applications": [
            app("大土方与装车", "Bulk earthmoving and loading", "中大规模土方要求挖掘、回转、提升和装车循环保持稳定产量。", "Medium-to-large earthmoving requires consistent output across digging, swing, boom lift and truck loading.", "循环时间、斗容、回转力矩和油耗", "Cycle time, bucket capacity, swing torque and fuel", "assets/ppt-insights/excavator-19-24t-01.webp", 154),
            app("基础与沟槽", "Foundations and trenching", "建筑基础和深沟需要长斗杆、精确控制、起吊和高密封驾驶环境。", "Building foundations and deep trenches require a long stick, precise control, lifting and a sealed operator environment.", "长斗杆、挖深、起吊和驾驶防护", "Long stick, digging depth, lift and cab protection", "assets/ppt-insights/excavator-19-24t-02.webp", 155),
            app("破碎与重型属具", "Breaker and heavy attachments", "标配 AUX1、AUX2、卸油、PTO 和合流能够提升机具匹配，但持续输出仍需温控和结构验证。", "Standard AUX1, AUX2, case drain, PTO and flow combining improve tool matching, while sustained output still requires thermal and structural validation.", "流量、压力、冷却、管路和防护", "Flow, pressure, cooling, plumbing and protection", "assets/ppt-insights/excavator-19-24t-03.webp", 157),
        ],
        "engineering": [
            row("型谱与工作范围", "Portfolio and working range", "基础参数与标杆总体接近。", "Base parameters are broadly close to the benchmark.", "缺少长斗杆，深沟作业受限。", "The lack of a long stick limits deep-trench work.", "配置长斗杆并同步建立起吊、稳定和运输数据。", "Offer a long stick and jointly establish lift, stability and transport data.", [156, 158]),
            row("液压与附件", "Hydraulics and attachments", "系统流量、AUX1/AUX2、卸油、PTO 和合流具备优势。", "System flow, AUX1/AUX2, case drain, PTO and flow combining are strengths.", "回转系统压力和铲斗挖掘力偏低。", "Swing-system pressure and bucket digging force are low.", "优化回转与铲斗机构，保持附件系统优势。", "Optimize swing and bucket linkage while preserving attachment-system strengths.", [157, 161]),
            row("操控与智能", "Controls and intelligence", "行走电控和阀芯优化改善了复合动作。", "Travel electronic control and spool optimization improved combined functions.", "停止上车动作时行走冲击较大，智能辅助和 360°环视不足。", "Travel shock occurs when upper functions stop, and guidance/360-degree visibility are insufficient.", "优化全电控标定并完成 360°与 2D/3D 功能验证。", "Optimize electrohydraulic calibration and validate 360-degree vision and 2D/3D functions.", [156, 160]),
            row("重载与安全", "Heavy duty and safety", "高密封增压驾驶室和全防护适合破碎作业。", "A sealed pressurized cab and full guarding support breaker work.", "起吊能力和回转装车效率弱于高配标杆。", "Lift and swing-loading productivity trail premium benchmarks.", "用载荷表、循环计时和冲击数据设定重载目标。", "Set heavy-duty targets using lift charts, cycle timing and shock data.", [157, 161]),
        ],
    },
    "excavator-24-28t": {
        "label": bi("24-28 吨级", "24-28 t class"), "model": "XCMG XE250U", "source_slides": [169, 187], "market_slide": 169,
        "concentration": bi("三个领先品牌合计约 75%", "Three leading brands hold about 75% combined"),
        "benchmark": bi("卡特 325/326 产品体系", "Caterpillar 325/326 portfolio"),
        "role": bi("24-28 吨市场同时存在常规尾和短尾，历史销量中短尾总体更好；常规尾仍承担大范围、重载和高效率任务。", "The 24-28 t market includes conventional- and short-tail machines, with short-tail historically stronger overall, while conventional-tail machines retain working-range, heavy-duty and productivity roles."),
        "transport": bi("XE250U 运输高度 3,210 mm、操作质量约 27,300 kg，资料采用三轴低底盘拖车并指出承载已接近参考边界，运输灵活性有限。", "XE250U transport height is 3,210 mm and operating mass about 27,300 kg. The source uses a three-axle lowboy and notes that payload is near the reference boundary, limiting transport flexibility."),
        "customers": bi("大型土方、道路桥梁、城市建设、基础与专业租赁客户", "Large earthmoving, road/bridge, urban construction, foundation and professional rental customers"),
        "jobs": bi("大土方装车、深沟、破碎、吊装、平地和重型附件", "Bulk loading, deep trenching, breaking, lifting, grading and heavy attachments"),
        "purchase": bi("生产率、挖掘力、回转、工作范围、附件系统、智能化和运输成本", "Productivity, digging forces, swing, working range, attachment system, intelligence and transport cost"),
        "applications": [
            app("重载土方", "Heavy earthmoving", "常规尾和较大工作范围支持大斗容装车，但最远端平地的斗杆控制需要保持稳定。", "Conventional-tail packaging and working range support larger-bucket loading, while far-reach grading requires stable stick control.", "斗容、回转、提升、平地微动和燃油", "Bucket capacity, swing, boom lift, grading modulation and fuel", "assets/ppt-insights/excavator-24-33t-01.webp", 172),
            app("深沟与基础", "Deep trench and foundations", "缺少长斗杆会限制深沟覆盖，负流量系统的微操控也会影响清底和边坡质量。", "Lack of a long stick limits deep-trench coverage, while negative-flow modulation affects cleanup and slope quality.", "长斗杆、挖深、微动和 2D/3D 引导", "Long stick, depth, fine control and 2D/3D guidance", "assets/ppt-insights/excavator-24-33t-02.webp", 173),
            app("装车与场地整形", "Loading and site shaping", "行走与上车复合、回转制动和平地协调直接影响循环和操作冲击。", "Travel/upper combined functions, swing braking and grading coordination directly affect cycle time and operating shock.", "复合动作、回转制动、平地和行走", "Combined functions, swing braking, grading and travel", "assets/ppt-insights/excavator-24-33t-03.webp", 174),
            app("破碎与重型机具", "Breaker and heavy tools", "AUX1、AUX2、卸油、PTO 和合流支持高流量机具，高密封驾驶室和防护构成安全基础。", "AUX1, AUX2, case drain, PTO and flow combining support high-flow tools; sealed cab and guarding provide the safety base.", "流量、压力、冷却、防护和结构耐久", "Flow, pressure, cooling, guarding and structural durability", "assets/ppt-insights/excavator-24-33t-04.webp", 175),
        ],
        "engineering": [
            row("整机与运输", "Machine and transport", "标配挖掘范围和恶劣路面行走驱动力具备优势。", "Standard working range and travel drive on difficult ground are strengths.", "三轴运输接近载荷边界。", "Three-axle transport is near the payload boundary.", "控制交付质量并评估附件拆分运输方案。", "Control delivered mass and evaluate split shipment of attachments.", [171, 177, 180]),
            row("液压与附件", "Hydraulics and attachments", "系统流量、AUX 和机具合流优于标杆。", "System flow, AUX and attachment-flow combining are strengths.", "负流量系统微操控和智能化不足。", "Negative-flow fine control and machine intelligence are insufficient.", "升级控制架构并保留附件液压优势。", "Upgrade the control architecture while retaining attachment-hydraulic strengths.", [178, 180]),
            row("重载性能", "Heavy-duty performance", "行走驱动力和系统流量适合泥泞及机具工况。", "Travel drive and system flow support soft-ground and tool applications.", "挖掘力和回转力矩弱于卡特，影响装车循环。", "Digging force and swing torque trail Caterpillar and affect loading cycles.", "按大斗容装车循环优化机构、回转和提升。", "Optimize linkage, swing and boom lift using large-bucket loading cycles.", [180]),
            row("智能与安全", "Intelligence and safety", "驾驶室密封与附件防护具备基础。", "Cab sealing and attachment protection provide a good base.", "缺少成熟 360°环视和高级智控。", "Mature 360-degree vision and advanced control functions are missing.", "完成 360°、2D/3D 和安全边界的整机验证。", "Complete machine validation of 360-degree vision, 2D/3D and safety boundaries.", [179]),
        ],
    },
    "excavator-24-28t-short-tail": {
        "label": bi("24-28 吨级短尾", "24-28 t short-tail class"), "model": "XCMG XE235UCR", "source_slides": [169, 187], "market_slide": 169,
        "concentration": bi("三个领先品牌合计约 75%", "Three leading brands hold about 75% combined"),
        "benchmark": bi("卡特 325/326 短尾体系", "Caterpillar 325/326 short-tail portfolio"),
        "role": bi("24-28 吨历史销量中短尾总体更好，核心客户在城市、居民区和道路项目中需要降低尾部碰撞风险，同时保留中挖生产率。", "Short-tail machines historically lead the 24-28 t range, serving urban, residential and road projects that need lower tail-strike risk without giving up mid-size productivity."),
        "transport": bi("XE235UCR 运输高度 3,181 mm、操作质量约 26,300 kg，资料显示可落在三轴低底盘参考承载范围内，但仍需超宽许可。", "XE235UCR transport height is 3,181 mm and operating mass about 26,300 kg, within the referenced three-axle lowboy payload while still requiring an oversize-width permit."),
        "customers": bi("城市建设、市政道路、居民区基础、公用工程和专业租赁客户", "Urban construction, municipal roads, residential foundations, utilities and professional rental fleets"),
        "jobs": bi("受限空间土方、沟槽、回填、吊装、破碎和多属具", "Confined earthmoving, trenching, backfill, lifting, breaking and multi-tool work"),
        "purchase": bi("短尾安全、运输、工作范围、复合动作、附件液压和驾驶视野", "Short-tail safety, transport, working range, combined functions, attachment hydraulics and visibility"),
        "applications": [
            app("城市受限土方", "Confined urban earthmoving", "短尾降低对车道和邻近建筑的侵入，最远端平地和回转制动决定施工质量与冲击。", "Short-tail packaging reduces lane and building intrusion; far-reach grading and swing braking determine quality and shock.", "尾部包络、视野、平地微动和回转控制", "Tail envelope, visibility, grading modulation and swing control", "assets/ppt-insights/excavator-24-33t-01.webp", 172),
            app("沟槽与公用工程", "Trenching and utilities", "在受限空间中仍需要长斗杆和精确清底，短尾不能替代工作范围和稳定验证。", "Confined work still requires a long stick and precise cleanup; short tail cannot replace working range and stability validation.", "长斗杆、挖深、起吊和稳定", "Long stick, depth, lift and stability", "assets/ppt-insights/excavator-24-33t-02.webp", 173),
            app("回填与整形", "Backfill and shaping", "行走与上车复合、回转制动和推土铲动作影响城市工程中的短循环效率。", "Travel/upper combined functions, swing braking and blade actions affect short-cycle productivity in urban work.", "复合动作、回转、行走和平地", "Combined functions, swing, travel and grading", "assets/ppt-insights/excavator-24-33t-03.webp", 174),
            app("破碎与防护", "Breaker work and protection", "高密封驾驶室、全防护和完整附件管路是城市拆改及重型机具作业的基础。", "A sealed cab, full guarding and complete attachment plumbing support urban demolition and heavy-tool work.", "驾驶防护、AUX、卸油、冷却和噪声", "Cab protection, AUX, case drain, cooling and noise", "assets/ppt-insights/excavator-24-33t-04.webp", 175),
        ],
        "engineering": [
            row("短尾与运输", "Short tail and transport", "短尾和较低交付质量更适合城市运输与施工。", "Short tail and lower delivered mass better suit urban transport and work.", "仍需超宽许可，附件组合会侵蚀载荷余量。", "Oversize-width permitting remains and attachment packages consume payload margin.", "建立城市运输包和附件拆装流程。", "Define an urban transport package and attachment removal process.", [169, 171]),
            row("工作范围", "Working range", "基础作业能力可覆盖主流城市任务。", "Base capability covers mainstream urban tasks.", "缺少长斗杆和超长工装，深沟覆盖不足。", "Long-stick and long-reach options are missing, limiting deep-trench coverage.", "优先配置长斗杆，超长工装按市场需求单独立项。", "Prioritize a long stick; treat long-reach equipment as a separate market-driven program.", [173, 176, 178]),
            row("操控与液压", "Controls and hydraulics", "AUX、PTO、合流和驾驶防护具备优势。", "AUX, PTO, flow combining and cab protection are strengths.", "复合动作停止冲击和微操控线性不足。", "Combined-function stop shock and fine-control linearity are insufficient.", "优化全电控回转与行走复合标定。", "Optimize electrohydraulic swing and travel-combination calibration.", [173, 175, 178]),
            row("智能安全", "Intelligent safety", "短尾本体降低碰撞风险。", "Short-tail packaging inherently reduces strike risk.", "360°环视和高级智能辅助尚未形成成熟配置。", "360-degree vision and advanced guidance are not yet mature equipment offers.", "将环视、电子围栏和 2D/3D 组成城市安全包验证。", "Validate an urban safety package combining surround view, e-fence and 2D/3D.", [179]),
        ],
    },
    "excavator-28-33t": {
        "label": bi("28-33 吨级", "28-33 t class"), "model": "XCMG XE300U", "source_slides": [170, 198], "market_slide": 170,
        "concentration": bi("三个领先品牌合计约 74%", "Three leading brands hold about 74% combined"),
        "benchmark": bi("卡特 330 产品体系", "Caterpillar 330 portfolio"),
        "role": bi("28-33 吨历史市场总量偏低、优先级低于相邻主销吨级，但现有负流量平台仍需要升级以维持基本竞争力。", "The 28-33 t market is historically smaller and lower priority than adjacent core classes, but the current negative-flow platform still requires modernization to remain competitive."),
        "transport": bi("XE300U 操作质量约 32,500 kg、运输高度 3,360 mm，资料采用四轴低底盘拖车并要求按超宽和轴荷复核。", "XE300U operating mass is about 32,500 kg and transport height 3,360 mm. The reference uses a four-axle lowboy with oversize-width and axle-load checks."),
        "customers": bi("中大型土方、道路桥梁、基础施工、采石辅助和专业租赁客户", "Medium/large earthmoving, road/bridge, foundation, quarry support and professional rental customers"),
        "jobs": bi("重载挖掘、装车、深沟、破碎、吊装和水利施工", "Heavy digging, truck loading, deep trenches, breaking, lifting and waterworks"),
        "purchase": bi("发动机与液压效率、挖掘力、起吊、回转、工装覆盖和智能控制", "Engine/hydraulic efficiency, digging force, lift, swing, work-equipment coverage and intelligent control"),
        "applications": [
            app("重载挖掘", "Heavy excavation", "系统流量和挖掘力支持重载，但发动机储备和控制架构决定持续生产率。", "System flow and digging force support heavy work, while engine reserve and control architecture determine sustained productivity.", "发动机功率、流量、挖掘力和冷却", "Engine power, flow, digging force and cooling", "assets/ppt-insights/excavator-24-33t-01.webp", 172),
            app("深沟与水利", "Deep trench and waterworks", "深沟和疏浚需要长斗杆或超长工装，缺少配置会直接限制投标与任务覆盖。", "Deep trenching and dredging require long-stick or long-reach equipment; missing options directly limit bids and task coverage.", "长斗杆、超长工装、稳定和作业范围", "Long stick, long reach, stability and working range", "assets/ppt-insights/excavator-24-33t-02.webp", 173),
            app("土方装车", "Earthmoving and loading", "挖掘、回转和提升循环需要动力、回转力矩和全电控协调。", "Dig, swing and boom cycles require coordinated power, swing torque and electrohydraulic control.", "循环、回转、提升、燃油和视野", "Cycle time, swing, boom lift, fuel and visibility", "assets/ppt-insights/excavator-24-33t-03.webp", 174),
            app("重型机具", "Heavy attachments", "大流量附件是徐工基础优势，持续作业仍需验证冷却、结构和回油。", "High-flow attachments are an XCMG foundation strength, while sustained work still requires cooling, structural and return-line validation.", "AUX、PTO、卸油、冷却和防护", "AUX, PTO, case drain, cooling and guarding", "assets/ppt-insights/excavator-24-33t-04.webp", 175),
        ],
        "engineering": [
            row("平台与动力", "Platform and power", "挖掘力和系统流量具备重载基础。", "Digging force and system flow provide a heavy-duty base.", "老旧负流量平台和发动机功率不足削弱持续效率。", "The aging negative-flow platform and limited engine power reduce sustained productivity.", "评估动力与全电控液压平台升级。", "Evaluate an engine and full electrohydraulic platform upgrade.", [170, 189, 197]),
            row("工作范围", "Working range", "基本参数与卡特接近。", "Base parameters are close to Caterpillar.", "挖掘深度略低，长斗杆与超长工装不足。", "Digging depth is lower and long-stick/long-reach coverage is insufficient.", "建立标准、长斗杆和长臂三层工装策略。", "Establish standard, long-stick and long-reach work-equipment tiers.", [176, 188, 189]),
            row("重载性能", "Heavy-duty performance", "挖掘力和流量优于历史标杆判断。", "Digging force and flow were rated favorably in the historical comparison.", "起吊能力和回转力矩略低。", "Lift and swing torque are slightly lower.", "通过配重、结构、回转和动臂提升联合优化。", "Jointly optimize counterweight, structure, swing and boom lift.", [191]),
            row("智能与安全", "Intelligence and safety", "附件液压具备扩展基础。", "Attachment hydraulics provide an expansion base.", "全电控智能化和高级安全功能落后。", "Full electrohydraulic intelligence and advanced safety features trail the benchmark.", "导入 2D/3D、称重、环视和电子围栏并量化收益。", "Introduce 2D/3D, payload, surround view and e-fence and quantify benefit.", [189, 190]),
        ],
    },
    "excavator-33-40t": {
        "label": bi("33-40 吨级", "33-40 t class"), "model": "XCMG XE360U", "source_slides": [199, 215], "market_slide": 199,
        "concentration": bi("三个领先品牌合计约 80%", "Three leading brands hold about 80% combined"),
        "benchmark": bi("卡特 336/340 产品体系", "Caterpillar 336/340 portfolio"),
        "role": bi("33-40 吨是中大挖热门规格，覆盖道路桥梁、城市建设、土方和初级石方；主流品牌同时布局常规尾、短尾和高低配。", "The 33-40 t range is a core large excavator class for roads, bridges, urban construction, earthmoving and entry-level rock work. Leading brands offer conventional-tail, short-tail and multiple performance tiers."),
        "transport": bi("资料采用四轴低底盘拖车，承载能力参考 38-45 t，并明确需要超宽、超高和超重许可。", "The reference uses a four-axle lowboy with 38-45 t payload capability and identifies oversize-width, height and weight permits."),
        "customers": bi("大型土方、道路桥梁、城市建设、采石场和设备租赁客户", "Large earthmoving, road/bridge, urban construction, quarry and rental customers"),
        "jobs": bi("大斗容装车、重载挖掘、道路基础、破碎、吊装和石方", "Large-bucket loading, heavy digging, road foundations, breaking, lifting and rock work"),
        "purchase": bi("吨产量、可靠性、运输、重型工装、回转效率、智能化和全生命周期成本", "Tons per hour, reliability, transport, heavy work equipment, swing productivity, intelligence and lifecycle cost"),
        "applications": [
            app("大型土方装车", "Large-scale loading", "高频装车依赖斗容、回转力矩、动臂提升和稳定性，型号高低配通常围绕这些能力分层。", "High-cycle loading depends on bucket capacity, swing torque, boom lift and stability, which often define model tiers.", "斗容、回转、提升、稳定和循环时间", "Bucket capacity, swing, boom lift, stability and cycle time", "assets/ppt-insights/excavator-33-40t-01.webp", 201),
            app("道路桥梁", "Road and bridge work", "基础开挖、回填、边坡和预制件吊装要求工作范围、起吊和智能平地。", "Foundation excavation, backfill, slopes and precast handling require working range, lift and intelligent grading.", "工作范围、起吊、2D/3D 和视野", "Working range, lift, 2D/3D and visibility", "assets/ppt-insights/excavator-33-40t-02.webp", 202),
            app("石方与破碎", "Rock and breaker work", "重型动臂、宽底盘、破碎管路和驾驶防护决定石方入门能力及持续耐久。", "Heavy boom, wide undercarriage, breaker plumbing and cab guarding determine entry-level rock capability and sustained durability.", "重型工装、底盘、AUX、冷却和防护", "Heavy work equipment, undercarriage, AUX, cooling and guarding", "assets/ppt-insights/excavator-33-40t-03.webp", 203),
            app("城市受限施工", "Confined urban work", "短尾机型补齐城市建设和人口密集区域需求，常规尾无法完全替代。", "A short-tail variant fills urban and dense-area demand that a conventional-tail machine cannot fully cover.", "短尾型谱、环视、运输和安全边界", "Short-tail portfolio, surround view, transport and safety envelope", "assets/ppt-insights/excavator-33-40t-04.webp", 204),
        ],
        "engineering": [
            row("型谱与运输", "Portfolio and transport", "常规尾 XE360U 覆盖主流土方。", "The conventional-tail XE360U covers mainstream earthmoving.", "缺少成熟短尾产品；运输高度和工作范围弱于标杆。", "A mature short-tail product is missing, and transport height/working range trail the benchmark.", "补齐短尾型谱并优化运输高度与标准工装。", "Complete the short-tail portfolio and optimize transport height and standard work equipment.", [199, 200, 205]),
            row("重载工装", "Heavy work equipment", "斗杆配置较全，挖掘力具备优势。", "Stick choices are broad and digging force is a strength.", "缺少重型动臂与更宽底盘组合。", "Heavy-boom and wider-undercarriage combinations are missing.", "形成土方、重载和石方工装包。", "Create earthmoving, heavy-duty and rock work-equipment packages.", [206, 208]),
            row("行走与回转", "Travel and swing", "基础行走和挖掘性能能够覆盖主流任务。", "Base travel and digging performance cover mainstream tasks.", "没有自动双速，回转力矩低于卡特。", "Automatic two-speed is missing and swing torque trails Caterpillar.", "升级自动双速并按装车循环优化回转与提升。", "Add automatic two-speed and optimize swing/boom lift using loading cycles.", [206, 208]),
            row("智能与安全", "Intelligence and safety", "现有平台可承载附件与辅助功能。", "The current platform can host attachment and assist functions.", "智控、环视和高级作业辅助不足。", "Machine intelligence, surround view and advanced work assists are insufficient.", "导入 2D/3D、称重、环视和电子围栏。", "Introduce 2D/3D, payload, surround view and e-fence.", [207]),
        ],
    },
    "excavator-40-60t": {
        "label": bi("40-60 吨级", "40-60 t class"), "model": "XCMG XE490U", "source_slides": [216, 232], "market_slide": 216,
        "concentration": bi("四个领先品牌合计超过 85%", "Four leading brands exceed 85% combined"),
        "benchmark": bi("卡特 349/352 产品体系", "Caterpillar 349/352 portfolio"),
        "role": bi("40-60 吨是高价值大挖规格，覆盖大型土方、道路桥梁、石方和入门级采矿；运输拆分能力已经成为产品竞争力的一部分。", "The 40-60 t range is a high-value large-excavator class for major earthmoving, roads/bridges, rock and entry-level mining. Transport disassembly capability is part of product competitiveness."),
        "transport": bi("资料采用五轴低底盘拖车，重型机具可能需要六轴；XE490U 高 4,400 mm、宽 3,592 mm，历史判断认为高度超过 4,267 mm 后需要引导车并显著增加成本。", "The reference uses a five-axle lowboy, potentially six axles with heavy tools. XE490U height is 4,400 mm and width 3,592 mm; the historical assessment indicates that height above 4,267 mm requires escort support and materially raises cost."),
        "customers": bi("大型土方、道路桥梁、采石场、矿山承包商和大型租赁客户", "Major earthmoving, road/bridge, quarry, mining contractors and large rental fleets"),
        "jobs": bi("大斗容装车、重载挖掘、破碎、石方、深沟、吊装和矿山辅助", "Large-bucket loading, heavy digging, breaking, rock, deep trenches, lifting and mine support"),
        "purchase": bi("吨产量、运输拆分、回转与起吊、重型工装、底盘、智能化、可靠性和服务能力", "Tons per hour, transport disassembly, swing and lift, heavy work equipment, undercarriage, intelligence, reliability and service capability"),
        "applications": [
            app("大型土方装车", "Large earthmoving and loading", "320 kW 动力和挖掘力支持大斗容作业，但回转、提升和稳定性决定最终循环产量。", "The 320 kW engine and digging force support large buckets, while swing, boom lift and stability determine final cycle output.", "斗容、回转、提升、稳定和燃油", "Bucket capacity, swing, boom lift, stability and fuel", "assets/ppt-insights/excavator-40-60t-01.webp", 218),
            app("石方与破碎", "Rock and breaker work", "高流量机具、PTO 和防护支持破碎作业，履带、结构、冷却和回转冲击决定耐久。", "High-flow tools, PTO and guarding support breaker work, while track system, structure, cooling and swing shock determine durability.", "AUX/PTO、履带、结构、冷却和防护", "AUX/PTO, tracks, structure, cooling and guarding", "assets/ppt-insights/excavator-40-60t-02.webp", 219),
            app("深沟与大型基础", "Deep trench and major foundations", "缺少长斗杆会限制深沟和大型基础任务，起吊与回转力矩决定管件和构件搬运效率。", "Lack of a long stick limits deep trenches and major foundations; lift and swing torque determine pipe and component handling.", "长斗杆、挖深、起吊、回转和稳定", "Long stick, depth, lift, swing and stability", "assets/ppt-insights/excavator-40-60t-03.webp", 220),
            app("采石与矿山辅助", "Quarry and mine support", "高负荷和长循环要求整机持续产量、可靠性、底盘耐久和快速服务。", "High loads and long cycles require sustained output, reliability, undercarriage durability and rapid service.", "持续功率、热平衡、底盘、远程管理和服务", "Sustained power, thermal balance, undercarriage, telematics and service", "assets/ppt-insights/excavator-40-60t-04.webp", 221),
        ],
        "engineering": [
            row("运输与底盘", "Transport and undercarriage", "整机宽度处于资料参考的引导车边界内。", "Machine width is within the referenced escort threshold.", "运输高度高、标配履带窄，缺少伸缩底盘和可拆配重。", "Transport height is high, standard tracks are narrow, and retractable undercarriage/removable counterweight are unavailable.", "以运输成本为目标开发伸缩底盘、可拆配重和履带选项。", "Develop retractable undercarriage, removable counterweight and track options against a transport-cost target.", [217, 222, 224, 231]),
            row("动力与附件", "Power and attachments", "320 kW 动力、挖掘力、AUX1/AUX2、PTO 和合流具备优势。", "The 320 kW engine, digging force, AUX1/AUX2, PTO and flow combining are strengths.", "长斗杆和更完整工装覆盖不足。", "Long-stick and broader work-equipment coverage are insufficient.", "建立土方、深沟、石方和矿山四类工装包。", "Create earthmoving, deep-trench, rock and mine work-equipment packages.", [223, 225, 231]),
            row("回转与起吊", "Swing and lifting", "重载挖掘力支持高负荷作业。", "Heavy digging force supports high-load work.", "回转力矩、起吊和动臂提升速度弱于卡特，影响装车。", "Swing torque, lift and boom speed trail Caterpillar and reduce loading productivity.", "按大斗容装车循环联合优化回转、提升、配重和稳定。", "Jointly optimize swing, boom lift, counterweight and stability using large-bucket loading cycles.", [225, 231]),
            row("智能与服务", "Intelligence and service", "机具系统为智能控制提供接口基础。", "The attachment system provides an integration base for intelligent control.", "智控、360°环视和运输辅助弱于高端标杆。", "Machine intelligence, 360-degree vision and transport aids trail the premium benchmark.", "导入作业辅助、环视、称重和远程健康管理。", "Introduce work assists, surround view, payload and remote health management.", [224, 231]),
        ],
    },
}


def chart_data(presentation: Presentation, slide_number: int | None) -> dict:
    if not slide_number:
        return {"years": [], "series": [], "status": []}
    slide = presentation.slides[slide_number - 1]
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        chart = shape.chart
        try:
            categories = [str(value) for value in chart.plots[0].categories]
        except (AttributeError, IndexError, TypeError):
            continue
        if len(categories) < 3:
            continue
        series = []
        for item in chart.series:
            try:
                values = [0 if value is None else int(round(float(value))) for value in item.values]
            except (AttributeError, TypeError, ValueError):
                values = []
            series.append({"brand": str(item.name or ""), "values": values})
        return {
            "years": categories,
            "series": series,
            "status": ["historical_actual", "historical_actual", "historical_estimate", "historical_forecast"],
        }
    return {"years": [], "series": [], "status": []}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    presentation = Presentation(args.source)
    records = []
    for slug, item in INSIGHTS.items():
        record = {"slug": slug, **item}
        record["market_chart"] = chart_data(presentation, item.get("market_slide"))
        records.append(record)

    payload = {
        "meta": {
            "source": args.source.name,
            "classification": "XCMG ARC INTERNAL",
            "scoring_boundary": "Market and field insight modules do not alter the Excel-based benchmark scores.",
            "market_time_basis": "2023-2024 historical actual; 2025 historical estimate; 2026 historical forecast pending actual backfill.",
        },
        "records": records,
    }
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"Wrote {len(records)} tonnage insight records to {args.output}")


if __name__ == "__main__":
    main()
