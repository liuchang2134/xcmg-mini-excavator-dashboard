import json
import re
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation

try:
    from tools.ppt_scope import OVERVIEW_SLIDES, display_title, slugs_for_slide
    from tools.ppt_source import source_pptx
except ModuleNotFoundError:
    from ppt_scope import OVERVIEW_SLIDES, display_title, slugs_for_slide
    from ppt_source import source_pptx


ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = source_pptx()
OUTPUT_PATH = ROOT / "data" / "ppt-insights" / "ppt-business-tables.json"

NAVIGATION_LABELS = {
    "销量分析",
    "适应性分析",
    "参配对比",
    "核心性能",
    "产品定位",
}

def normalize_text(value):
    value = (value or "").replace("\r", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return normalize_text(shape.text)


def slide_texts(slide):
    candidates = []
    for shape in slide.shapes:
        value = shape_text(shape)
        if not value:
            continue
        candidates.append(
            {
                "text": value,
                "top": int(shape.top),
                "left": int(shape.left),
                "width": int(shape.width),
                "height": int(shape.height),
            }
        )
    return sorted(candidates, key=lambda item: (item["top"], item["left"]))


def choose_slide_title(items, slide_number):
    texts = [item["text"] for item in items]
    preferred = [
        value
        for value in texts
        if any(
            token in value
            for token in (
                "核心规格产品",
                "产品线竞争力",
                "可售型谱",
                "未来三年",
                "市场洞察资源",
                "宏观环境分析",
                "行业分析",
                "各品牌占有率",
            )
        )
        and len(value) <= 160
    ]
    if preferred:
        return preferred[0]
    for value in texts:
        if len(value) <= 120 and not value.startswith(("二、", "三、")):
            return value
    return f"第 {slide_number} 页业务表格"


def trim_matrix(rows):
    rows = [[normalize_text(cell) for cell in row] for row in rows]
    while rows and not any(rows[-1]):
        rows.pop()
    if not rows:
        return []
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    while width and not any(row[width - 1] for row in rows):
        width -= 1
    return [row[:width] for row in rows]


def is_navigation_table(rows):
    if len(rows) != 1 or len(rows[0]) != 5:
        return False
    cells = {
        re.sub(r"\s+", "", cell)
        for cell in rows[0]
        if normalize_text(cell)
    }
    matches = sum(
        any(label in cell for label in NAVIGATION_LABELS)
        for cell in cells
    )
    return matches >= 3


def classify_role(title, slide_number):
    title = normalize_text(title)
    if slide_number == 8:
        return "environment"
    if slide_number in {10, 12, 13, 14}:
        return "competition"
    if "产品定位" in title or "占有率目标" in title:
        return "positioning"
    if "销量" in title or "占有率" in title or "市场份额" in title:
        return "market"
    if "适应性分析" in title:
        return "application"
    if "参数及配置" in title:
        return "product_comparison"
    if "型谱" in title:
        return "portfolio"
    if "提升举措" in title:
        return "roadmap"
    if slide_number == 245:
        return "sales_plan"
    if slide_number == 247:
        return "evidence_system"
    return "field_evaluation"


def table_heading(title, table_index, table_count):
    title = normalize_text(title)
    if table_count <= 1:
        return title
    return f"{title} · 表 {table_index}"


def extract():
    deck = Presentation(PPTX_PATH)
    records = []
    excluded = Counter()

    for slide_number, slide in enumerate(deck.slides, start=1):
        if slide_number == 246:
            excluded["personnel_table"] += sum(
                1 for shape in slide.shapes if getattr(shape, "has_table", False)
            )
            continue

        texts = slide_texts(slide)
        source_title = choose_slide_title(texts, slide_number)
        title = display_title(slide_number, source_title)
        table_shapes = [
            shape for shape in slide.shapes if getattr(shape, "has_table", False)
        ]
        accepted = []
        for shape in table_shapes:
            rows = trim_matrix(
                [[cell.text for cell in row.cells] for row in shape.table.rows]
            )
            if not rows or not any(any(row) for row in rows):
                excluded["empty_table"] += 1
                continue
            if is_navigation_table(rows):
                excluded["navigation_table"] += 1
                continue
            accepted.append(rows)

        for table_index, rows in enumerate(accepted, start=1):
            width = max((len(row) for row in rows), default=0)
            record_id = f"s{slide_number:03d}-t{table_index:02d}"
            records.append(
                {
                    "id": record_id,
                    "slide": slide_number,
                    "table_index": table_index,
                    "title": table_heading(title, table_index, len(accepted)),
                    "source_title_zh": source_title,
                    "title_correction": (
                        {
                            "status": "corrected_source_heading",
                            "reason": "Source heading conflicts with the table tonnage, models and chapter context.",
                        }
                        if title != source_title
                        else None
                    ),
                    "role": classify_role(title, slide_number),
                    "rows": len(rows),
                    "columns": width,
                    "matrix_zh": rows,
                    "overview": slide_number in OVERVIEW_SLIDES,
                    "slugs": slugs_for_slide(slide_number),
                    "historical_source": True,
                }
            )

    by_slug = defaultdict(list)
    overview = []
    for record in records:
        if record["overview"]:
            overview.append(record["id"])
        for slug in record["slugs"]:
            by_slug[slug].append(record["id"])

    payload = {
        "meta": {
            "source_file": PPTX_PATH.name,
            "slide_count": len(deck.slides),
            "classification": "XCMG ARC INTERNAL",
            "scope": "All business-meaningful source tables; repeated navigation and personnel tables excluded.",
            "language_note": "Original Chinese table cells are retained verbatim to protect data fidelity.",
        },
        "summary": {
            "included_table_count": len(records),
            "overview_table_count": len(overview),
            "tonnage_table_count": sum(len(items) for items in by_slug.values()),
            "excluded": dict(excluded),
            "roles": dict(Counter(record["role"] for record in records)),
        },
        "overview": overview,
        "by_slug": dict(by_slug),
        "records": records,
    }
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    OUTPUT_PATH.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return payload


if __name__ == "__main__":
    result = extract()
    print(json.dumps(result["summary"], ensure_ascii=False, indent=2))
