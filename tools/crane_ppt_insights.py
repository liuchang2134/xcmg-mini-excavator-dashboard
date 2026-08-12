from __future__ import annotations

import hashlib
import json
import os
import re
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SOURCE = (
    ROOT
    / "data"
    / "source-presentations"
    / "XCMG_North_America_Crane_Market_Insight_2025-07-01_V13_source.pptx"
)
OUTPUT_DIR = ROOT / "data" / "crane-ppt-insights"
ASSET_DIR = ROOT / "assets" / "crane-ppt-source"
SOURCE_DATE = "2025-07-01"
MIN_IMAGE_SLIDE_AREA_RATIO = 0.007
MIN_IMAGE_PIXEL_AREA = 25_000

CLASS_SECTION_SLIDES = {
    "RT-60t": {
        "market-insight": [81],
        "job-applications": list(range(82, 86)),
        "engineering-insight": [86, *range(88, 94)],
        "product-positioning": [94],
    },
    "RT-75t": {
        "market-insight": [81],
        "job-applications": list(range(82, 86)),
        "engineering-insight": list(range(87, 94)),
        "product-positioning": [95],
    },
    "RT-100t": {
        "market-insight": [81],
        "job-applications": list(range(96, 100)),
        "engineering-insight": list(range(100, 107)),
        "product-positioning": [107],
    },
    "RT-130t": {
        "market-insight": [81],
        "job-applications": list(range(108, 115)),
        "engineering-insight": list(range(115, 119)),
        "product-positioning": [119],
    },
    "RT-160t": {
        "market-insight": [81, 132],
        "job-applications": [],
        "engineering-insight": [],
        "product-positioning": [147, 152],
    },
    "AT-150t": {
        "market-insight": [58],
        "job-applications": list(range(59, 63)),
        "engineering-insight": list(range(63, 68)),
        "product-positioning": [68],
    },
}

CLASS_SLIDES = {
    class_id: list(
        dict.fromkeys(
            slide
            for section_slides in sections.values()
            for slide in section_slides
        )
    )
    for class_id, sections in CLASS_SECTION_SLIDES.items()
}

CLASS_META = {
    "RT-60t": {
        "title_zh": "60吨级越野轮胎起重机",
        "title_en": "60 t Rough-Terrain Crane",
        "models": ["XCR60_U", "XCR75_U"],
        "source_scope": "current-at-source-date",
    },
    "RT-75t": {
        "title_zh": "75吨级越野轮胎起重机",
        "title_en": "75 t Rough-Terrain Crane",
        "models": ["XCR60_U", "XCR75_U"],
        "source_scope": "current-at-source-date",
    },
    "RT-100t": {
        "title_zh": "100吨级越野轮胎起重机",
        "title_en": "100 t Rough-Terrain Crane",
        "models": ["XCR100_U"],
        "source_scope": "current-at-source-date",
    },
    "RT-130t": {
        "title_zh": "130吨级越野轮胎起重机",
        "title_en": "130 t Rough-Terrain Crane",
        "models": ["XCR130_U"],
        "source_scope": "current-at-source-date",
    },
    "RT-160t": {
        "title_zh": "160-165吨级越野轮胎起重机",
        "title_en": "160-165 US ton Rough-Terrain Crane",
        "models": ["XCR165U"],
        "source_scope": "plan",
    },
    "AT-150t": {
        "title_zh": "150吨级全地面起重机",
        "title_en": "150 t All-Terrain Crane",
        "models": ["XCA150_U"],
        "source_scope": "current-at-source-date",
    },
}

PLAN_SLIDES = set(range(143, 164)) | {26, 42, 57, 68, 80, 94, 95, 107, 119, 129}


def _clean_text(value: Any) -> str:
    text = str(value or "").replace("\x0b", "\n").replace("\xa0", " ")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _geometry(shape: Any) -> dict[str, int]:
    return {
        "left": int(getattr(shape, "left", 0) or 0),
        "top": int(getattr(shape, "top", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
    }


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return _clean_text(shape.text_frame.text)


def _extract_text_blocks(shapes: list[Any]) -> list[dict[str, Any]]:
    blocks = []
    for shape in shapes:
        text = _shape_text(shape)
        if not text:
            continue
        blocks.append({"text": text, **_geometry(shape)})
    return sorted(blocks, key=lambda item: (item["top"], item["left"]))


def _extract_tables(shapes: list[Any]) -> list[dict[str, Any]]:
    tables = []
    for index, shape in enumerate(shapes, 1):
        if not getattr(shape, "has_table", False):
            continue
        rows = [
            [_clean_text(cell.text) for cell in row.cells]
            for row in shape.table.rows
        ]
        if not any(any(cell for cell in row) for row in rows):
            continue
        tables.append({"id": f"table-{index:02d}", "rows": rows, **_geometry(shape)})
    return tables


def _chart_categories(chart: Any) -> list[str]:
    try:
        return [_clean_text(category.label) for category in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError):
        return []


def _chart_xml_values(series: Any, tag: str) -> list[float | str | None]:
    """Read scatter and bubble values not exposed by python-pptx's public API."""
    try:
        nodes = series._element.xpath(f"./c:{tag}")
    except (AttributeError, TypeError):
        return []
    if not nodes:
        return []
    values: list[float | str | None] = []
    for raw in nodes[0].xpath(".//c:pt/c:v/text()"):
        try:
            values.append(float(raw))
        except (TypeError, ValueError):
            values.append(_clean_text(raw))
    return values


def _extract_charts(shapes: list[Any]) -> list[dict[str, Any]]:
    charts = []
    for index, shape in enumerate(shapes, 1):
        if not getattr(shape, "has_chart", False):
            continue
        chart = shape.chart
        series = []
        for item in chart.series:
            values = []
            for value in item.values:
                if value is None:
                    values.append(None)
                else:
                    try:
                        values.append(float(value))
                    except (TypeError, ValueError):
                        values.append(_clean_text(value))
            series.append(
                {
                    "name": _clean_text(item.name),
                    "values": values,
                    "x_values": _chart_xml_values(item, "xVal"),
                    "y_values": _chart_xml_values(item, "yVal"),
                    "bubble_sizes": _chart_xml_values(item, "bubbleSize"),
                }
            )
        charts.append(
            {
                "id": f"chart-{index:02d}",
                "chart_type": str(chart.chart_type),
                "categories": _chart_categories(chart),
                "series": series,
                **_geometry(shape),
            }
        )
    return charts


def _browser_image(blob: bytes, ext: str, target: Path) -> Path | None:
    ext = ext.lower().strip(".")
    if ext in {"jpg", "jpeg", "png", "gif", "webp"}:
        output = target.with_suffix(".jpg" if ext == "jpeg" else f".{ext}")
        output.write_bytes(blob)
        return output
    try:
        with Image.open(BytesIO(blob)) as image:
            image.load()
            output = target.with_suffix(".webp")
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGBA" if "transparency" in image.info else "RGB")
            image.save(output, "WEBP", lossless=True, quality=95)
            return output
    except (OSError, ValueError):
        return None


def _extract_images(
    shapes: list[Any],
    slide_number: int,
    slide_width: int,
    slide_height: int,
) -> list[str]:
    images = []
    slide_area = max(slide_width * slide_height, 1)
    for index, shape in enumerate(shapes, 1):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        geometry = _geometry(shape)
        shape_ratio = geometry["width"] * geometry["height"] / slide_area
        pixel_width, pixel_height = shape.image.size
        pixel_area = pixel_width * pixel_height
        # The presentation uses compact field photos and diagnostic close-ups.
        # They are business evidence, even when each occupies only a small part
        # of the original slide, so retain them while filtering tiny ornaments.
        if shape_ratio < MIN_IMAGE_SLIDE_AREA_RATIO or pixel_area < MIN_IMAGE_PIXEL_AREA:
            continue
        blob = shape.image.blob
        digest = hashlib.sha256(blob).hexdigest()[:10]
        stem = ASSET_DIR / f"s{slide_number:03d}-image-{index:02d}-{digest}"
        output = _browser_image(blob, shape.image.ext, stem)
        if output is None:
            continue
        images.append(output.relative_to(ROOT).as_posix())
    return images


def _section_for_slide(slide: int) -> str:
    if slide <= 2:
        return "macro"
    if slide == 3:
        return "market-volume"
    if slide <= 6:
        return "competition"
    if slide <= 9:
        return "regional-demand"
    if slide <= 57:
        return "boom-truck"
    if slide <= 80:
        return "all-terrain"
    if slide <= 119:
        return "rough-terrain"
    if slide <= 129:
        return "crawler"
    if slide <= 142:
        return "portfolio"
    if slide <= 152:
        return "roadmap"
    return "go-to-market"


def _status_for_slide(slide: int) -> str:
    if slide in PLAN_SLIDES:
        return "plan"
    if slide in {3, 4, 5, 7, 8, 9, 10, 58, 81, 120}:
        return "historical"
    return "current-at-source-date"


def _title_from_blocks(slide: int, blocks: list[dict[str, Any]]) -> str:
    ignored = {"二、大区产业板块洞察", "二、大区产业板块洞察—北美"}
    for block in blocks:
        candidates = [line.strip() for line in block["text"].splitlines() if line.strip()]
        for candidate in candidates:
            if candidate not in ignored and len(candidate) > 3:
                return candidate[:180]
    return f"Slide {slide}"


def _source_path() -> Path:
    override = os.environ.get("XCMG_CRANE_PPT", "").strip()
    return Path(override) if override else DEFAULT_SOURCE


def _write_json(name: str, payload: Any) -> None:
    path = OUTPUT_DIR / name
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def build_crane_ppt_insights(source: Path | None = None) -> dict[str, Any]:
    source_path = Path(source) if source else _source_path()
    if not source_path.exists():
        raise FileNotFoundError(f"Crane insight presentation not found: {source_path}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    for stale in ASSET_DIR.glob("s*-image-*"):
        stale.unlink()

    presentation = Presentation(str(source_path))
    records = []
    evidence = []
    source_hash = hashlib.sha256(source_path.read_bytes()).hexdigest()
    for slide_number, slide in enumerate(presentation.slides, 1):
        shapes = list(_iter_shapes(slide.shapes))
        blocks = _extract_text_blocks(shapes)
        tables = _extract_tables(shapes)
        charts = _extract_charts(shapes)
        images = _extract_images(
            shapes,
            slide_number,
            presentation.slide_width,
            presentation.slide_height,
        )
        classes = [name for name, numbers in CLASS_SLIDES.items() if slide_number in numbers]
        record = {
            "slide": slide_number,
            "title": _title_from_blocks(slide_number, blocks),
            "section": _section_for_slide(slide_number),
            "status": _status_for_slide(slide_number),
            "validation_status": "plan-not-verified" if slide_number in PLAN_SLIDES else "source-recorded",
            "source_date": SOURCE_DATE,
            "source_type": "pptx",
            "classes": classes,
            "text_blocks": blocks,
            "tables": tables,
            "charts": charts,
            "images": images,
        }
        records.append(record)
        evidence.append(
            {
                "id": f"crane-slide-{slide_number:03d}",
                "source_slide": slide_number,
                "title": record["title"],
                "section": record["section"],
                "status": record["status"],
                "validation_status": record["validation_status"],
                "source_date": SOURCE_DATE,
                "classes": classes,
                "assets": images,
            }
        )

    segment_map = {
        name: {**CLASS_META[name], "slides": slides}
        for name, slides in CLASS_SLIDES.items()
    }
    _write_json("slides.json", records)
    _write_json("segment-map.json", segment_map)
    _write_json("evidence.json", evidence)
    _write_json(
        "source.json",
        {
            "file_name": source_path.name,
            "sha256": source_hash,
            "slide_count": len(records),
            "source_date": SOURCE_DATE,
            "generated_assets": sum(len(item["images"]) for item in records),
            "native_tables": sum(len(item["tables"]) for item in records),
            "native_charts": sum(len(item["charts"]) for item in records),
        },
    )
    return {
        "slide_count": len(records),
        "asset_count": sum(len(item["images"]) for item in records),
        "table_count": sum(len(item["tables"]) for item in records),
        "chart_count": sum(len(item["charts"]) for item in records),
    }


if __name__ == "__main__":
    result = build_crane_ppt_insights()
    print(json.dumps(result, ensure_ascii=False, indent=2))
