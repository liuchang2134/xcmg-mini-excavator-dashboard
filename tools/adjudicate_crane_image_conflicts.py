from __future__ import annotations

import hashlib
import json
import re
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
DECISIONS_PATH = ROOT / "data" / "crane-ppt-insights" / "image-ownership.json"
REUSE_REPORT_PATH = ROOT / "crane-source-reuse-report.json"
SLIDES_PATH = ROOT / "data" / "crane-ppt-insights" / "slides.json"
SOURCE_DIR = ROOT / "assets" / "crane-ppt-source"
PRESENTATION_DIR = ROOT / "data" / "source-presentations"


def _iter_pictures(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pictures(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _presentation_image_pages(presentation: Path) -> dict[str, list[int]]:
    pages: dict[str, set[int]] = defaultdict(set)
    deck = Presentation(str(presentation))
    for slide_number, slide in enumerate(deck.slides, 1):
        for picture in _iter_pictures(slide.shapes):
            digest = hashlib.md5(picture.image.blob).hexdigest()
            pages[digest].add(slide_number)
    return {digest: sorted(values) for digest, values in pages.items()}


def _presentation_media_reuse(presentation: Path) -> list[dict[str, Any]]:
    with zipfile.ZipFile(presentation) as archive:
        media = {
            name: hashlib.md5(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        }
        references: dict[str, set[int]] = defaultdict(set)
        media_names: dict[str, set[str]] = defaultdict(set)
        for name in archive.namelist():
            match = re.fullmatch(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", name)
            if not match:
                continue
            slide_number = int(match.group(1))
            relationships = archive.read(name).decode("utf-8", "ignore")
            for target in re.findall(r'Target="\.\./media/([^"]+)"', relationships):
                media_path = f"ppt/media/{target}"
                digest = media.get(media_path)
                if not digest:
                    continue
                references[digest].add(slide_number)
                media_names[digest].add(target)
    groups = [
        {
            "content_hash": digest,
            "slides": sorted(slides),
            "media": sorted(media_names[digest]),
        }
        for digest, slides in references.items()
        if len(slides) > 1
    ]
    return sorted(groups, key=lambda item: (item["slides"], item["content_hash"]))


def _slide_number(path: str) -> int:
    match = re.search(r"(?:^|/)s(\d+)-", path.replace("\\", "/"))
    if not match:
        raise ValueError(f"Cannot read slide number from {path}")
    return int(match.group(1))


def _source_asset(path: str) -> str:
    candidate = ROOT / path
    if "crane-ppt-source" in candidate.as_posix() and candidate.exists():
        return candidate.relative_to(ROOT).as_posix()
    stem = candidate.stem
    matches = sorted(SOURCE_DIR.glob(f"{stem}.*"))
    if len(matches) != 1:
        raise ValueError(f"Expected one source asset for {path}; found {len(matches)}")
    return matches[0].relative_to(ROOT).as_posix()


def adjudicate(presentation: Path | None = None) -> dict[str, Any]:
    records = json.loads(SLIDES_PATH.read_text(encoding="utf-8"))
    if presentation is None:
        presentations = sorted(PRESENTATION_DIR.glob("*.pptx"))
        if len(presentations) != 1:
            raise FileNotFoundError(
                f"Expected one source presentation in {PRESENTATION_DIR}; found {len(presentations)}"
            )
        presentation = presentations[0]
    presentation = Path(presentation)
    ppt_pages = _presentation_image_pages(presentation)
    assets: dict[str, list[int]] = defaultdict(list)
    for record in records:
        for asset in record.get("images") or []:
            assets[str(asset)].append(int(record["slide"]))
    decisions = []
    asset_decisions: dict[str, dict[str, Any]] = {}

    reused_assets = []
    for asset, referenced_pages in assets.items():
        digest = hashlib.md5((ROOT / asset).read_bytes()).hexdigest()
        source_pages = ppt_pages.get(digest, [])
        if len(source_pages) > 1:
            reused_assets.append((asset, digest, source_pages, referenced_pages))

    for index, (asset, digest, source_pages, referenced_pages) in enumerate(
        sorted(reused_assets, key=lambda item: (_slide_number(item[0]), item[0])), 1
    ):
        pages_zh = "、".join(str(page) for page in source_pages)
        pages_en = ", ".join(str(page) for page in source_pages)
        decision = {
            "id": f"reuse-{index:02d}",
            "content_hash": digest,
            "decision": "SOURCE_REUSE",
            "confidence": 1.0,
            "source_slides": source_pages,
            "reason_zh": f"PPT 原始文件在第 {pages_zh} 页复用了完全相同的图片二进制，不能据此判断唯一机型或区域。",
            "reason_en": f"The source presentation reuses the exact same image binary on slides {pages_en}; it cannot establish a unique model or region.",
            "source_assets": [asset],
            "referenced_slides": sorted(set(referenced_pages)),
        }
        decisions.append(decision)
        asset_decisions[asset] = {
            "decision_id": decision["id"],
            "decision": decision["decision"],
            "source_slides": source_pages,
        }

    reuse_groups = _presentation_media_reuse(presentation)
    REUSE_REPORT_PATH.write_text(
        json.dumps(
            {
                "source_presentation": presentation.resolve().relative_to(ROOT.resolve()).as_posix(),
                "method": "Exact MD5 reuse across PPT media relationships; includes all media references, not only rendered picture shapes.",
                "reuse_group_count": len(reuse_groups),
                "groups": reuse_groups,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    payload = {
        "source_presentation": presentation.resolve().relative_to(ROOT.resolve()).as_posix(),
        "method": "Exact binary-image reuse verified inside the source PPTX. Reader-facing captions remain business-focused; slide and reuse traceability is retained in page data attributes.",
        "decision_count": len(decisions),
        "decisions": decisions,
        "assets": asset_decisions,
    }
    DECISIONS_PATH.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return payload


if __name__ == "__main__":
    result = adjudicate()
    print(
        f"Auto-adjudicated {result['decision_count']} conflict groups as verified source reuse; "
        f"wrote {DECISIONS_PATH.relative_to(ROOT)}"
    )
