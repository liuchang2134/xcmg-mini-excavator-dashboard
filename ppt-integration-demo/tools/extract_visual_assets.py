"""Extract reusable pictures, tables and chart data from selected PPT slides.

The output is an internal build artifact for the web prototype. It deliberately
works at shape level so the page never needs to display a full-slide screenshot.
"""

from __future__ import annotations

import argparse
import io
import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def crop_picture(shape) -> Image.Image:
    image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
    width, height = image.size
    left = round(width * shape.crop_left)
    top = round(height * shape.crop_top)
    right = round(width * (1 - shape.crop_right))
    bottom = round(height * (1 - shape.crop_bottom))
    if right <= left or bottom <= top:
        return image
    return image.crop((left, top, right, bottom))


def shape_box(shape) -> dict[str, float]:
    return {
        "left_in": round(shape.left / 914400, 3),
        "top_in": round(shape.top / 914400, 3),
        "width_in": round(shape.width / 914400, 3),
        "height_in": round(shape.height / 914400, 3),
    }


def table_data(shape) -> list[list[str]]:
    return [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]


def chart_data(shape) -> dict:
    chart = shape.chart
    categories = []
    try:
        categories = [str(value) for value in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError):
        pass
    series = []
    for item in chart.series:
        values = []
        try:
            values = [None if value is None else float(value) for value in item.values]
        except (AttributeError, TypeError, ValueError):
            pass
        series.append({"name": str(item.name or ""), "values": values})
    return {"chart_type": str(chart.chart_type), "categories": categories, "series": series}


def extract(source: Path, output_dir: Path, manifest_path: Path, slides: list[int]) -> None:
    presentation = Presentation(source)
    output_dir.mkdir(parents=True, exist_ok=True)
    manifest = {
        "source": source.name,
        "scope": slides,
        "rule": "Shape-level extraction only; full-slide images are excluded from the web UI.",
        "slides": [],
    }

    for slide_number in slides:
        slide = presentation.slides[slide_number - 1]
        record = {"slide": slide_number, "pictures": [], "tables": [], "charts": [], "text": []}
        picture_index = 0
        table_index = 0
        chart_index = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_index += 1
                rendered = crop_picture(shape)
                filename = f"s{slide_number:03d}-photo-{picture_index:02d}.jpg"
                rendered.save(output_dir / filename, "JPEG", quality=91, optimize=True)
                record["pictures"].append(
                    {
                        "file": filename,
                        "shape_name": shape.name,
                        "pixel_width": rendered.width,
                        "pixel_height": rendered.height,
                        **shape_box(shape),
                    }
                )
            elif getattr(shape, "has_table", False):
                table_index += 1
                record["tables"].append(
                    {"id": f"s{slide_number:03d}-table-{table_index:02d}", "rows": table_data(shape), **shape_box(shape)}
                )
            elif getattr(shape, "has_chart", False):
                chart_index += 1
                record["charts"].append(
                    {"id": f"s{slide_number:03d}-chart-{chart_index:02d}", **chart_data(shape), **shape_box(shape)}
                )

            if getattr(shape, "has_text_frame", False):
                value = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                if value:
                    record["text"].append({"value": value, **shape_box(shape)})

        manifest["slides"].append(record)

    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    pictures = sorted(output_dir.glob("*.jpg"))
    if pictures:
        tiles: list[Image.Image] = []
        for path in pictures:
            image = Image.open(path).convert("RGB")
            image.thumbnail((340, 210))
            tile = Image.new("RGB", (360, 245), "white")
            tile.paste(image, ((360 - image.width) // 2, 26))
            ImageDraw.Draw(tile).text((8, 7), path.name, fill="black")
            tiles.append(tile)
        columns = 4
        rows = (len(tiles) + columns - 1) // columns
        contact_sheet = Image.new("RGB", (columns * 360, rows * 245), (225, 230, 235))
        for index, tile in enumerate(tiles):
            contact_sheet.paste(tile, ((index % columns) * 360, (index // columns) * 245))
        contact_sheet.save(output_dir / "contact-sheet.jpg", "JPEG", quality=88)
    print(
        json.dumps(
            {
                "slides": len(manifest["slides"]),
                "pictures": sum(len(item["pictures"]) for item in manifest["slides"]),
                "tables": sum(len(item["tables"]) for item in manifest["slides"]),
                "charts": sum(len(item["charts"]) for item in manifest["slides"]),
            },
            ensure_ascii=False,
        )
    )


def parse_slides(value: str) -> list[int]:
    result: set[int] = set()
    for part in value.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = (int(number) for number in part.split("-", 1))
            result.update(range(start, end + 1))
        else:
            result.add(int(part))
    return sorted(result)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output_dir", type=Path)
    parser.add_argument("manifest", type=Path)
    parser.add_argument("--slides", default="3-15,48-68")
    args = parser.parse_args()
    extract(args.source, args.output_dir, args.manifest, parse_slides(args.slides))


if __name__ == "__main__":
    main()
